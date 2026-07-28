#!/usr/bin/env python3
"""Generate Sahil's 29 July 2026 quarterly business review artifacts."""

from __future__ import annotations

import argparse
from collections import Counter, defaultdict
from datetime import date, datetime
from pathlib import Path
from typing import Any

from openpyxl import Workbook, load_workbook
from openpyxl.styles import Alignment, Font, PatternFill
from openpyxl.utils.datetime import from_excel
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


AS_OF = date(2026, 7, 27)
REVIEW_DATE = date(2026, 7, 29)
REPORTING_OWNERS = {"Sahil Patni", "Kasturi Rangan"}

NAVY = "0B132B"
NAVY_2 = "111C3A"
TEAL = "00B8A9"
BLUE = "3A86FF"
ORANGE = "FF9F1C"
RED = "FF5A5F"
WHITE = "F7FAFC"
MUTED = "A8B2C7"
GRID = "2A3658"
PALE = "DDE6F3"
GREEN = "3DDC97"

REFERENCES = [
    (
        "R1",
        "International Federation of Robotics, World Robotics 2025 press release",
        "https://ifr.org/ifr-press-releases/news/global-robot-demand-in-factories-doubles-over-10-years",
    ),
    (
        "R2",
        "Rockwell Automation, 2025 State of Smart Manufacturing — APAC findings",
        "https://www.rockwellautomation.com/en-sg/company/news/press-releases/apac-sosm-2025.html",
    ),
    (
        "R3",
        "Interact Analysis, Machine Vision return-to-growth forecast, June 2025",
        "https://interactanalysis.com/return-to-growth-forecast-for-machine-vision-in-2025-despite-us-tariffs/",
    ),
    (
        "R4",
        "Altera, Robotics Solutions Stack",
        "https://www.altera.com/fpga-solutions/robotics-solutions-stack",
    ),
    (
        "R5",
        "Altera, Agilex 5 FPGA and SoC FPGA overview",
        "https://www.altera.com/products/fpga/agilex/5",
    ),
    (
        "R6",
        "AMD, Kria KR260 robotics platform and Kria Robotics Stack",
        "https://www.amd.com/en/products/system-on-modules/kria/k26/robotics.html",
    ),
    (
        "R7",
        "Altera, Video and Vision Processing Suite",
        "https://www.altera.com/products/ip/po-3150/video-and-vision-processing-suite",
    ),
    (
        "R8",
        "Altera, FPGA AI Suite",
        "https://www.altera.com/products/development-tools/fpga-ai-suite",
    ),
    (
        "R9",
        "AMD, Kria KV260 Vision AI Starter Kit",
        "https://www.amd.com/en/products/system-on-modules/kria/k26/kv260-vision-starter-kit.html",
    ),
    (
        "R10",
        "AMD, Versal AI Edge Series",
        "https://www.amd.com/en/products/adaptive-socs-and-fpgas/versal/ai-edge-series.html",
    ),
    (
        "R11",
        "Altera, Industrial solutions",
        "https://www.altera.com/fpga-solutions/industrial",
    ),
    (
        "R12",
        "Altera, Agilex 5 SoC HPS features including three 2.5G TSN Ethernet MACs",
        "https://docs.altera.com/r/docs/762191/current/agilextm-5-fpgas-and-socs-device-overview/additional-features-for-agilextm-5-socs",
    ),
    (
        "R13",
        "AMD, Industrial Networking solutions",
        "https://www.amd.com/en/solutions/industrial/industrial-networking.html",
    ),
    (
        "R14",
        "Altera, Agilex 5 Functional Safety",
        "https://docs.altera.com/api/khub/documents/xhgkkZ1PZaLEHFnNiUlbNA/content",
    ),
    (
        "R15",
        "AMD, Functional Safety",
        "https://www.amd.com/en/products/adaptive-socs-and-fpgas/technologies/functional-safety.html",
    ),
    (
        "R16",
        "Altera, Agilex 3 FPGA and SoC FPGA overview",
        "https://www.altera.com/products/fpga/agilex/3",
    ),
    (
        "R17",
        "AMD, Spartan UltraScale+ FPGA overview",
        "https://www.amd.com/en/products/adaptive-socs-and-fpgas/fpga/spartan-ultrascale-plus.html",
    ),
    (
        "R18",
        "AMD, Vitis unified software platform",
        "https://www.amd.com/en/products/software/adaptive-socs-and-fpgas/vitis.html",
    ),
    (
        "R19",
        "Altera, Video Solutions Stack",
        "https://www.altera.com/fpga-solutions/video-solutions-stack",
    ),
    (
        "R20",
        "Altera, Sensor Interfaces including Holoscan Sensor Bridge",
        "https://www.altera.com/fpga-solutions/sensory-interfaces",
    ),
]

PRIMARY_ACCOUNTS = ["Ciena", "Juniper / HPE", "Philips", "GE HealthCare", "Outdu"]
OTHER_STRATEGIC_ACCOUNTS = [
    "Siemens",
    "Boeing",
    "Honeywell Aerospace",
    "Schneider",
    "Emerson",
]

ALLOCATION_ROSTER = [
    ("ADA", "Bangalore", "ADG", "08 Feb 2025", "Amol"),
    ("ADE", "Bangalore", "ADG", "08 Feb 2025", "Amol"),
    ("DEAL", "Dehradun", "ADG", "08 Feb 2025", "Mayank"),
    ("C-DOT", "Across India", "COMM", "08 Jan 2026", "Mayank / Amol"),
    ("BEL", "Bangalore", "ADG", "08 Feb 2025", "Sumanth / Prasad"),
    ("LRDE", "Bangalore", "ADG", "08 Feb 2025", "Sumanth / Prasad"),
    ("ECON", "Chennai", "OTHERS", "08 Feb 2025", "Akash"),
    ("BEL", "Chennai", "ADG", "08 Feb 2025", "Akash"),
    ("IRDE", "Dehradun", "ADG", "08 Feb 2025", "TBH / Prasad"),
    ("BEL", "Ghaziabad", "ADG", "08 Feb 2025", "TBH / Prasad"),
    ("BEL", "Machilipatnam", "ADG", "08 Feb 2025", "Sumanth"),
    ("Outdu", "Bangalore", "OTHERS", "08 Jan 2026", "TBD"),
]

SUPPORT_ISSUES = [
    (
        "Quartus Pro 24.1 SignalTap compilation slowdown",
        "C-DOT",
        date(2026, 5, 22),
        "Tools / AI",
        "Reproduce with an example design and SignalTap; provide a cleaner debug path.",
        "Reproduction required",
    ),
    (
        "LVDS SerDes + optical transceiver integration for TSE",
        "Macnica / external",
        date(2026, 5, 29),
        "Networking",
        "100 Mbps optical fit requires hardware adaptation; TSE LVDS path is 1G-oriented.",
        "Feasibility validation",
    ),
    (
        "Non-PTP drops with PTP residence-time update in LL 40GbE MAC",
        "Ceragon",
        date(2026, 5, 19),
        "Networking",
        "Trigger isolated to residence-time update control, narrowing MAC-level debug.",
        "Debug narrowed",
    ),
    (
        "MAX 10 DK-DEV-10M08E144-B GPIO feasibility",
        "Siemens",
        date(2026, 6, 25),
        "Device / I/O",
        "44-signal fit is feasible subject to pin, bank-voltage, shared-circuit and level-shift checks.",
        "Guidance provided",
    ),
    (
        "Agilex 7 F-Tile TX MAC segmented-interface behavior",
        "C-DOT",
        date(2026, 6, 18),
        "Networking",
        "Clarified fixed latency, idle behavior and mid-frame valid-drop error handling.",
        "Resolved",
    ),
    (
        "Cyclone V PCIe MSI MsiReq_o not asserting",
        "HPE",
        date(2026, 6, 29),
        "PCIe / Configuration",
        "Use MSI memory-write transaction flow rather than expecting MSIReq_o assertion.",
        "Guidance provided",
    ),
    (
        "MAX 10 I/O utilization for 10M02SCE144A7G",
        "Honeywell Aerospace",
        date(2026, 6, 26),
        "Device / I/O",
        "No fixed derating percentage; validate package and electrical constraints in Quartus.",
        "Guidance provided",
    ),
    (
        "PTA 26.1 LPDDR5 frequency cap versus datasheet",
        "Vicharak",
        date(2026, 7, 28),
        "Tools / AI",
        "Use 1866 MHz in Platform Designer and full Quartus power analysis pending PTA fix.",
        "Workaround + escalation",
    ),
    (
        "FPGA AI Suite HL-JTAG memory and address-map issue",
        "Pantherun",
        date(2026, 7, 15),
        "Tools / AI",
        "Use fabric DDR or expand the accessible bridge/span for F2SDRAM.",
        "Guidance provided",
    ),
    (
        "Agilex 7 custom-board bring-up and power sequencing",
        "BEL / Macnica",
        date(2026, 7, 27),
        "Board / Schematic",
        "Validate rail sequencing, pre-bias, waveforms, JTAG visibility and board captures first.",
        "Bring-up investigation",
    ),
    (
        "Agilex 3/5 schematic corrections before card release",
        "Qbit Labs / Arrow",
        date(2026, 7, 15),
        "Board / Schematic",
        "Correct GTS bank power grouping before AI-Modem and TSN card release.",
        "Release blocker found",
    ),
    (
        "ATUM Nano Ethernet failure after BSP regeneration",
        "Macnica / BEL Chennai",
        date(2026, 7, 3),
        "Networking",
        "Reapply required manual driver edits after BSP regeneration, then rebuild.",
        "Fix identified",
    ),
    (
        "Agilex 7 core/HSSI rail topology for 3U VPX switch",
        "BEL / Macnica",
        date(2026, 7, 7),
        "Board / Schematic",
        "Do not collapse regulators without current-sharing, transient and margin validation.",
        "Validation required",
    ),
    (
        "Dual-image Cyclone V GX flash update partitioning",
        "Juniper / HPE",
        date(2026, 7, 24),
        "PCIe / Configuration",
        "Need a partition-safe primary-image update method; generated RPD spans full flash.",
        "Method required",
    ),
    (
        "Cyclone V PCIe MSI interrupt generation",
        "Juniper / HPE",
        date(2026, 7, 14),
        "PCIe / Configuration",
        "Implement MSI through the memory-write flow rather than MsiReq_o.",
        "Guidance provided",
    ),
    (
        "SDR schematic-freeze pin and constraint blockers",
        "Macnica SDR program",
        date(2026, 7, 16),
        "Board / Schematic",
        "Close OPN mismatch, HPS UART, electrical constraints and pin handling before release.",
        "Release blockers found",
    ),
]


def rgb(hex_value: str) -> RGBColor:
    return RGBColor.from_string(hex_value)


def parse_date(value: Any) -> date | None:
    if isinstance(value, datetime):
        return value.date()
    if isinstance(value, date):
        return value
    if isinstance(value, (int, float)):
        return from_excel(value).date()
    if isinstance(value, str):
        for pattern in ("%m/%d/%Y", "%Y-%m-%d"):
            try:
                return datetime.strptime(value, pattern).date()
            except ValueError:
                pass
    return None


def fmt_value(value: float, decimals: int = 2) -> str:
    """Format in source units; the workbook does not identify a currency."""
    if abs(value) >= 1_000_000:
        return f"{value / 1_000_000:.{decimals}f}M"
    if abs(value) >= 1_000:
        return f"{value / 1_000:.0f}K"
    return f"{value:,.0f}"


def fmt_date(value: date | None) -> str:
    return value.strftime("%d %b %Y") if value else "—"


def normalize_probability(value: Any) -> float:
    probability = float(value or 0)
    return probability * 100 if 0 < probability <= 1 else probability


def read_pipeline(source: Path) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
    workbook = load_workbook(source, data_only=True)
    worksheet = workbook.active
    values = list(worksheet.iter_rows(values_only=True))
    headers = values[0]
    line_items = [
        dict(zip(headers, row))
        for row in values[1:]
        if row[2] in REPORTING_OWNERS
    ]

    opportunities: dict[str, dict[str, Any]] = {}
    for row in line_items:
        opportunity_id = row["Opportunity ID"]
        if opportunity_id not in opportunities:
            opportunities[opportunity_id] = {
                **row,
                "Peak Value": 0.0,
                "Products": [],
            }
        opportunities[opportunity_id]["Peak Value"] += float(row["Peak Value"] or 0)
        opportunities[opportunity_id]["Products"].append(row["Product Name"])

    result = list(opportunities.values())
    for row in result:
        row["Design Win Date Parsed"] = parse_date(row["Design Win Date"])
        row["Created Date Parsed"] = parse_date(row["Created Date"])
        row["Production Start Date Parsed"] = parse_date(row["Production Start Date"])
        probability = normalize_probability(row["Probability (%)"])
        row["Probability (%)"] = probability
        row["Weighted Value"] = row["Peak Value"] * probability / 100

    result.sort(key=lambda item: item["Peak Value"], reverse=True)
    return line_items, result


def read_open_pipeline_buckets(
    sources: dict[str, Path],
    canonical_opportunities: list[dict[str, Any]],
) -> dict[str, list[dict[str, Any]]]:
    """Map user-provided planning exports to canonical opportunity records."""
    canonical = {row["Opportunity ID"]: row for row in canonical_opportunities}
    buckets: dict[str, list[dict[str, Any]]] = {}
    seen_ids: set[str] = set()

    for bucket, source in sources.items():
        workbook = load_workbook(source, data_only=True)
        worksheet = workbook["Sahil-Open Pipeline"]
        values = list(worksheet.iter_rows(values_only=True))
        headers = values[0]
        rows = [
            dict(zip(headers, row))
            for row in values[1:]
            if row[2] in REPORTING_OWNERS
        ]
        ids = list(dict.fromkeys(row["Opportunity ID"] for row in rows))
        unknown_ids = [opportunity_id for opportunity_id in ids if opportunity_id not in canonical]
        if unknown_ids:
            raise ValueError(
                f"{source} contains reporting-scope opportunities absent from the canonical report: {unknown_ids}"
            )
        duplicate_ids = seen_ids.intersection(ids)
        if duplicate_ids:
            raise ValueError(f"Open-pipeline buckets overlap: {sorted(duplicate_ids)}")
        seen_ids.update(ids)
        bucket_rows = []
        for opportunity_id in ids:
            row = dict(canonical[opportunity_id])
            row["Open Pipeline Bucket"] = bucket
            bucket_rows.append(row)
        bucket_rows.sort(key=lambda item: item["Peak Value"], reverse=True)
        buckets[bucket] = bucket_rows

    return buckets


def add_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    size: float = 16,
    color: str = WHITE,
    bold: bool = False,
    font: str = "Aptos",
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.04,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.font.name = font
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = rgb(color)
    return shape


def add_rich_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    lines: list[tuple[str, str, float, bool]],
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.05)
    frame.margin_bottom = Inches(0.05)
    for index, (text, color, size, bold) in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = text
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = rgb(color)
        paragraph.space_after = Pt(5)
    return shape


def add_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = NAVY_2,
    line: str = GRID,
    radius: bool = True,
):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(1)
    return shape


def add_header(slide, title: str, section: str, number: int):
    add_text(slide, 0.55, 0.25, 2.7, 0.25, section.upper(), size=9, color=TEAL, bold=True)
    add_text(slide, 0.55, 0.52, 11.9, 0.48, title, size=25, bold=True)
    add_text(
        slide,
        11.95,
        0.26,
        0.75,
        0.25,
        f"{number:02d}",
        size=10,
        color=MUTED,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.08), Inches(12.15), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(GRID)
    line.line.fill.background()
    add_text(
        slide,
        0.55,
        7.18,
        12.15,
        0.18,
        "Prepared 27 Jul 2026  •  Source: Sahil Report.xlsx  •  Values are in source units (currency not encoded)",
        size=7.5,
        color=MUTED,
    )


def add_citation(slide, text: str):
    add_text(slide, 0.62, 6.88, 12.0, 0.18, text, size=7.2, color=MUTED)


def new_slide(prs: Presentation, title: str, section: str, number: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = rgb(NAVY)
    add_header(slide, title, section, number)
    return slide


def add_stat_card(slide, x: float, y: float, w: float, label: str, value: str, note: str, accent: str = TEAL):
    add_box(slide, x, y, w, 1.12)
    add_text(slide, x + 0.18, y + 0.13, w - 0.36, 0.22, label.upper(), size=8.5, color=MUTED, bold=True)
    add_text(slide, x + 0.18, y + 0.38, w - 0.36, 0.38, value, size=23, color=accent, bold=True)
    add_text(slide, x + 0.18, y + 0.82, w - 0.36, 0.18, note, size=8, color=PALE)


def add_table(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    headers: list[str],
    rows: list[list[str]],
    widths: list[float],
    *,
    font_size: float = 8.5,
    header_size: float = 8,
    row_fills: list[str] | None = None,
):
    table = slide.shapes.add_table(
        len(rows) + 1,
        len(headers),
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    ).table
    for index, width in enumerate(widths):
        table.columns[index].width = Inches(width)
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(BLUE)
        cell.margin_left = Inches(0.06)
        cell.margin_right = Inches(0.04)
        cell.margin_top = Inches(0.03)
        cell.margin_bottom = Inches(0.02)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(header_size)
        paragraph.font.bold = True
        paragraph.font.color.rgb = rgb(WHITE)
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    for row_index, row in enumerate(rows, start=1):
        fill = row_fills[row_index - 1] if row_fills else (NAVY_2 if row_index % 2 else "152446")
        for col, value in enumerate(row):
            cell = table.cell(row_index, col)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(fill)
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.025)
            cell.margin_bottom = Inches(0.02)
            frame = cell.text_frame
            frame.word_wrap = True
            frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            paragraph = frame.paragraphs[0]
            paragraph.font.name = "Aptos"
            paragraph.font.size = Pt(font_size)
            paragraph.font.color.rgb = rgb(WHITE)
    return table


def suggested_gate(opportunity: dict[str, Any]) -> str:
    stage = opportunity["Opportunity Stage"]
    if stage == "Identify":
        return "Qualify use case, sponsor, budget and 0→25% gate"
    if stage == "Define":
        return "Freeze architecture/BOM and evaluation plan"
    if stage == "Develop":
        return "Close validation plan and design-freeze evidence"
    if stage == "Design":
        return "Secure DW evidence and production handoff"
    return "Confirm next-stage exit criteria"


def customer_name(value: str) -> str:
    replacements = {
        "EMERSON INNOVATION CENTER - PUNE": "Emerson",
        "GE HEALTHCARE": "GE HealthCare",
        "HONEYWELL AEROSPACE INDIA PRIVATE LIMITED": "Honeywell Aerospace",
        "HONEYWELL TECHNOLOGY SOLUTIONS": "Honeywell Tech.",
        "JUNIPER NETWORKS": "Juniper",
        "Outdu Mediatech Private Limited": "Outdu",
        "Philips India Limited - Innovation Center": "Philips",
        "CIENA INDIA PVT LTD": "Ciena",
        "Boeing International Corporation India Pvt. Ltd.": "Boeing",
    }
    return replacements.get(value, value)


def account_in_group(group: str, account_name: str) -> bool:
    value = account_name.upper()
    checks = {
        "Ciena": ("CIENA",),
        "Juniper / HPE": ("JUNIPER", "HPE"),
        "Philips": ("PHILIPS",),
        "GE HealthCare": ("GE HEALTHCARE",),
        "Outdu": ("OUTDU",),
        "Siemens": ("SIEMENS",),
        "Boeing": ("BOEING",),
        "Honeywell Aerospace": ("HONEYWELL AEROSPACE",),
        "Schneider": ("SCHNEIDER",),
        "Emerson": ("EMERSON",),
    }
    return any(token in value for token in checks[group])


def issue_in_group(group: str, issue_customer: str) -> bool:
    value = issue_customer.upper()
    checks = {
        "Ciena": ("CIENA",),
        "Juniper / HPE": ("JUNIPER", "HPE"),
        "Philips": ("PHILIPS",),
        "GE HealthCare": ("GE HEALTHCARE", "GEHC"),
        "Outdu": ("OUTDU",),
        "Siemens": ("SIEMENS",),
        "Boeing": ("BOEING",),
        "Honeywell Aerospace": ("HONEYWELL",),
        "Schneider": ("SCHNEIDER",),
        "Emerson": ("EMERSON",),
    }
    return any(token in value for token in checks[group])


def build_presentation(
    opportunities: list[dict[str, Any]],
    open_buckets: dict[str, list[dict[str, Any]]],
    output: Path,
):
    total = sum(row["Peak Value"] for row in opportunities)
    open_2026 = open_buckets["2026"]
    open_2027 = open_buckets["2027"]
    open_opportunities = open_2026 + open_2027
    open_ids = {row["Opportunity ID"] for row in open_opportunities}
    outside_open_opportunities = [
        row for row in opportunities if row["Opportunity ID"] not in open_ids
    ]
    identify_opportunities = [
        row
        for row in outside_open_opportunities
        if row["Opportunity Stage"] == "Identify"
    ]
    other_outside_opportunities = [
        row
        for row in outside_open_opportunities
        if row["Opportunity Stage"] != "Identify"
    ]
    ranked_opportunities = sorted(
        open_opportunities, key=lambda row: row["Peak Value"], reverse=True
    )
    qualified_total = sum(row["Peak Value"] for row in open_opportunities)
    weighted = sum(row["Weighted Value"] for row in open_opportunities)
    customers = sorted({row["Account Name"] for row in open_opportunities})
    dw_2026 = [
        row
        for row in open_opportunities
        if row["Design Win Date Parsed"] and row["Design Win Date Parsed"].year == 2026
    ]
    q4_2026 = [
        row
        for row in open_opportunities
        if row["Design Win Date Parsed"]
        and date(2026, 10, 1) <= row["Design Win Date Parsed"] <= date(2026, 12, 31)
    ]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 — Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(NAVY)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = rgb(TEAL)
    accent.line.fill.background()
    add_text(slide, 0.8, 0.72, 4.0, 0.28, "QUARTERLY BUSINESS REVIEW", size=10, color=TEAL, bold=True)
    add_text(slide, 0.8, 1.22, 8.6, 1.32, "Own the next gate.\nBuild the 2027 engine.", size=35, bold=True)
    add_text(slide, 0.8, 2.78, 7.9, 0.42, "FAE plan review  •  Altera + Arrow + Macnica", size=18, color=PALE)
    add_text(slide, 0.8, 3.25, 7.6, 0.32, "Wednesday, 29 July 2026  |  12:00–13:00", size=13, color=MUTED)
    add_box(slide, 9.25, 0.75, 3.25, 5.55, fill=NAVY_2, line=GRID)
    add_text(slide, 9.62, 1.18, 2.5, 0.28, "PIPELINE SNAPSHOT", size=9, color=MUTED, bold=True)
    add_text(slide, 9.62, 1.65, 2.5, 0.58, fmt_value(qualified_total), size=31, color=TEAL, bold=True)
    add_text(slide, 9.62, 2.18, 2.5, 0.26, "qualified open value", size=10, color=PALE)
    add_text(slide, 9.62, 2.72, 2.5, 0.58, str(len(open_opportunities)), size=31, color=BLUE, bold=True)
    add_text(slide, 9.62, 3.25, 2.5, 0.26, "qualified open opportunities", size=10, color=PALE)
    add_text(slide, 9.62, 3.78, 2.5, 0.58, str(len(customers)), size=31, color=ORANGE, bold=True)
    add_text(slide, 9.62, 4.31, 2.5, 0.26, "direct-account customers", size=10, color=PALE)
    add_text(slide, 9.62, 5.02, 2.5, 0.68, "Target + actual\ninputs required", size=15, color=RED, bold=True)
    add_text(slide, 0.8, 6.78, 8.3, 0.24, "Prepared from master + 2026/2027 open-pipeline exports  •  Values shown without assumed currency", size=8, color=MUTED)

    # 2 — Executive snapshot
    slide = new_slide(
        prs,
        f"Executive snapshot: {fmt_value(qualified_total)} qualified open; {fmt_value(total - qualified_total)} outside exports",
        "01 / Position",
        2,
    )
    card_width = 2.82
    add_stat_card(slide, 0.55, 1.35, card_width, "Qualified open", fmt_value(qualified_total), f"{len(open_opportunities)} opportunities", TEAL)
    add_stat_card(slide, 3.55, 1.35, card_width, "Weighted pipeline", fmt_value(weighted), "probability × peak", BLUE)
    add_stat_card(slide, 6.55, 1.35, card_width, "2026 open export", fmt_value(sum(x["Peak Value"] for x in open_2026)), f"{len(open_2026)} opportunities", ORANGE)
    add_stat_card(slide, 9.55, 1.35, card_width, "2027 open export", fmt_value(sum(x["Peak Value"] for x in open_2027)), f"{len(open_2027)} opportunities", GREEN)
    add_box(slide, 0.55, 2.78, 7.55, 3.85)
    add_rich_text(
        slide,
        0.78,
        3.05,
        7.05,
        3.25,
        [
            ("WHAT THE DATA SAYS", TEAL, 10, True),
            (f"• {fmt_value(sum(x['Peak Value'] for x in outside_open_opportunities))} across {len(outside_open_opportunities)} plays is outside both exports: {fmt_value(sum(x['Peak Value'] for x in identify_opportunities))} Identify + {fmt_value(sum(x['Peak Value'] for x in other_outside_opportunities))} Define.", WHITE, 14, False),
            (f"• Outdu AI is {fmt_value(4_125_000)} ({4_125_000 / qualified_total:.0%} of qualified open value) but remains Define.", WHITE, 15, False),
            (f"• Q4 2026 DWIN-dated pipeline is {fmt_value(sum(x['Peak Value'] for x in q4_2026))}, yet weighted value is only {fmt_value(sum(x['Weighted Value'] for x in q4_2026))}.", WHITE, 15, False),
            ("• Bucket membership comes from the uploaded files; it is not derived from Design Win Date.", WHITE, 15, False),
        ],
    )
    add_box(slide, 8.35, 2.78, 4.35, 3.85, fill="14213D", line=ORANGE)
    add_rich_text(
        slide,
        8.62,
        3.05,
        3.78,
        3.20,
        [
            ("DECISIONS FOR THE REVIEW", ORANGE, 10, True),
            ("1  Confirm 2026 target and YTD actual.", WHITE, 15, True),
            (f"2  Agree stage-exit evidence for {len(q4_2026)} Q4 DWINs.", WHITE, 15, True),
            ("3  Assign DFAE support to top conversion plays.", WHITE, 15, True),
            ("4  Select a distributor plan or confirm direct-only coverage.", WHITE, 15, True),
        ],
    )

    # 3 — Target vs achievement
    slide = new_slide(prs, "2026 target vs current achievement: source inputs are incomplete", "01 / Position", 3)
    add_text(slide, 0.55, 1.28, 12.0, 0.32, "Do not use open opportunity value as booked achievement or target attainment.", size=13, color=ORANGE, bold=True)
    labels = [
        ("2026 ANNUAL TARGET", "[ INPUT REQUIRED ]", "Sales target is not in the source file"),
        ("YTD ACTUAL ACHIEVEMENT", "[ INPUT REQUIRED ]", "Bookings/revenue/DWIN actuals are not in the source file"),
        ("GAP TO TARGET", "[ CALCULATE AFTER INPUT ]", "Target − YTD actual"),
    ]
    for index, (label, value, note) in enumerate(labels):
        x = 0.55 + index * 4.08
        add_box(slide, x, 1.85, 3.82, 1.62, fill="14213D", line=RED if index < 2 else GRID)
        add_text(slide, x + 0.2, 2.05, 3.42, 0.22, label, size=8.5, color=MUTED, bold=True)
        add_text(slide, x + 0.2, 2.40, 3.42, 0.34, value, size=17, color=RED if index < 2 else ORANGE, bold=True)
        add_text(slide, x + 0.2, 2.94, 3.42, 0.28, note, size=8, color=PALE)
    add_box(slide, 0.55, 3.82, 12.15, 2.45)
    add_text(slide, 0.8, 4.10, 11.65, 0.25, "AVAILABLE LEADING INDICATORS", size=9.5, color=TEAL, bold=True)
    indicators = [
        ("Qualified open", f"{len(open_opportunities)} / {fmt_value(qualified_total)}"),
        ("Weighted pipeline", fmt_value(weighted)),
        ("2026 export", f"{len(open_2026)} / {fmt_value(sum(x['Peak Value'] for x in open_2026))}"),
        ("2027 export", f"{len(open_2027)} / {fmt_value(sum(x['Peak Value'] for x in open_2027))}"),
    ]
    for index, (label, value) in enumerate(indicators):
        x = 0.82 + index * 2.9
        add_text(slide, x, 4.58, 2.45, 0.22, label, size=9, color=MUTED, bold=True)
        add_text(slide, x, 4.90, 2.45, 0.38, value, size=19, color=WHITE, bold=True)
    add_text(slide, 0.82, 5.72, 11.2, 0.30, "Bring to review: annual target, YTD revenue/bookings, achieved DWIN count/value, and definition of “achievement.”", size=11.5, color=ORANGE, bold=True)

    # 4 — Funnel
    slide = new_slide(prs, "Opportunity funnel: separate qualified open from 0% identification", "02 / Funnel", 4)
    stage_order = ["Identify", "Define", "Develop", "Design"]
    stage_colors = [RED, ORANGE, BLUE, GREEN]
    stage_rows = []
    for stage, stage_color in zip(stage_order, stage_colors):
        rows = [row for row in opportunities if row["Opportunity Stage"] == stage]
        value = sum(row["Peak Value"] for row in rows)
        stage_rows.append((stage, len(rows), value, sum(row["Weighted Value"] for row in rows), stage_color))
    max_value = max(row[2] for row in stage_rows)
    for index, (stage, count, value, stage_weighted, stage_color) in enumerate(stage_rows):
        y = 1.45 + index * 1.05
        add_text(slide, 0.62, y + 0.08, 1.15, 0.22, stage, size=12, bold=True)
        add_text(slide, 1.58, y + 0.08, 0.65, 0.22, f"{count} opps", size=9, color=MUTED)
        add_box(slide, 2.35, y, 6.35, 0.48, fill="152446", line="152446", radius=False)
        bar_width = max(0.12, 6.35 * value / max_value)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.35), Inches(y), Inches(bar_width), Inches(0.48))
        bar.fill.solid()
        bar.fill.fore_color.rgb = rgb(stage_color)
        bar.line.fill.background()
        add_text(slide, 8.88, y + 0.03, 1.15, 0.24, fmt_value(value), size=13, bold=True, align=PP_ALIGN.RIGHT)
        add_text(slide, 10.18, y + 0.03, 1.52, 0.24, f"weighted {fmt_value(stage_weighted)}", size=9, color=MUTED)
    add_box(slide, 0.55, 5.72, 12.15, 0.75, fill="14213D", line=GRID)
    top_three = sum(row["Peak Value"] for row in ranked_opportunities[:3])
    add_text(slide, 0.82, 5.94, 3.35, 0.26, f"Qualified top 3 = {top_three / qualified_total:.0%}", size=13, color=ORANGE, bold=True)
    add_text(slide, 4.35, 5.94, 3.75, 0.26, f"Open exports = {fmt_value(qualified_total)}", size=13, color=TEAL, bold=True)
    add_text(slide, 8.35, 5.94, 3.85, 0.26, f"Identify outside exports = {fmt_value(total - qualified_total)}", size=13, color=RED, bold=True)

    # 5 — 2026 closure plan
    slide = new_slide(prs, f"2026 DWIN-dated pipeline: convert {len(dw_2026)} named plays", "03 / Closure", 5)
    ordered_dw = sorted(dw_2026, key=lambda row: row["Design Win Date Parsed"])
    rows = [
        [
            customer_name(row["Account Name"]),
            row["Opportunity Name"],
            row["Opportunity Stage"],
            f"{row['Probability (%)']:.0f}%",
            fmt_date(row["Design Win Date Parsed"]),
            fmt_value(row["Peak Value"]),
            suggested_gate(row),
        ]
        for row in ordered_dw
    ]
    add_table(
        slide,
        0.55,
        1.38,
        12.15,
        3.55,
        ["Customer", "Opportunity", "Stage", "Prob.", "DWIN", "Peak", "Proposed next gate"],
        rows,
        [1.15, 2.42, 0.76, 0.62, 1.03, 0.77, 5.40],
        font_size=8.4,
    )
    add_box(slide, 0.55, 5.22, 12.15, 1.18, fill="14213D", line=ORANGE)
    add_text(slide, 0.82, 5.46, 2.65, 0.22, "REVIEW COMMITMENT", size=9, color=ORANGE, bold=True)
    add_text(
        slide,
        3.05,
        5.39,
        9.15,
        0.62,
        "For each play: customer decision date • technical gap • required sample/tool/IP • named Sales + FAE + DFAE owners • evidence for the next stage.",
        size=13,
        color=WHITE,
        bold=True,
    )

    # 6/7 — Top opportunities
    for slide_number, chunk_start in ((6, 0), (7, 7)):
        chunk_end = 7 if chunk_start == 0 else 15
        chunk = ranked_opportunities[chunk_start:chunk_end]
        title = "Top qualified opportunities 1–7: protect the concentration" if chunk_start == 0 else "Top qualified opportunities 8–15: build breadth"
        slide = new_slide(prs, title, "04 / Top opportunities", slide_number)
        table_rows = []
        fills = []
        for rank, row in enumerate(chunk, start=chunk_start + 1):
            table_rows.append(
                [
                    str(rank),
                    customer_name(row["Account Name"]),
                    row["Opportunity Name"],
                    row["Opportunity Stage"],
                    f"{row['Probability (%)']:.0f}%",
                    fmt_date(row["Design Win Date Parsed"]),
                    fmt_value(row["Peak Value"]),
                    suggested_gate(row),
                ]
            )
            fills.append("14213D" if rank <= 3 else (NAVY_2 if rank % 2 else "152446"))
        add_table(
            slide,
            0.55,
            1.35,
            12.15,
            4.92,
            ["#", "Customer", "Opportunity", "Stage", "Prob.", "DWIN", "Peak", "Proposed next gate"],
            table_rows,
            [0.35, 1.25, 2.52, 0.72, 0.60, 1.00, 0.74, 4.97],
            font_size=8.1,
            row_fills=fills,
        )
        if chunk_start == 0:
            add_text(slide, 0.65, 6.43, 11.9, 0.22, f"Qualified ranks 1–3 contribute {top_three / qualified_total:.1%} of qualified open value; inspect weekly.", size=10.5, color=ORANGE, bold=True)
        else:
            add_text(slide, 0.65, 6.43, 11.9, 0.22, f"{len(open_opportunities)} qualified open opportunities exist; three lower-value plays sit below this Top 15.", size=10.5, color=ORANGE, bold=True)

    # 8 — Customer portfolio
    slide = new_slide(prs, "Qualified customer portfolio: Outdu and Ciena lead open value", "05 / Accounts", 8)
    grouped: dict[str, list[dict[str, Any]]] = defaultdict(list)
    for row in open_opportunities:
        grouped[row["Account Name"]].append(row)
    customer_rows = sorted(
        [
            (
                customer_name(account),
                len(rows),
                sum(item["Peak Value"] for item in rows),
                sum(item["Weighted Value"] for item in rows),
            )
            for account, rows in grouped.items()
        ],
        key=lambda row: row[2],
        reverse=True,
    )
    max_customer = max(row[2] for row in customer_rows)
    for index, (account, count, value, customer_weighted) in enumerate(customer_rows):
        y = 1.40 + index * 0.69
        add_text(slide, 0.62, y + 0.03, 1.65, 0.22, account, size=9.5, bold=True)
        add_box(slide, 2.35, y, 5.25, 0.34, fill="152446", line="152446", radius=False)
        bar_width = max(0.09, 5.25 * value / max_customer)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.35), Inches(y), Inches(bar_width), Inches(0.34))
        bar.fill.solid()
        bar.fill.fore_color.rgb = rgb(TEAL if index < 2 else BLUE)
        bar.line.fill.background()
        add_text(slide, 7.78, y - 0.01, 0.85, 0.22, fmt_value(value), size=10.5, bold=True, align=PP_ALIGN.RIGHT)
        add_text(slide, 8.76, y - 0.01, 1.10, 0.22, f"W {fmt_value(customer_weighted)}", size=8.5, color=MUTED)
        add_text(slide, 10.06, y - 0.01, 0.85, 0.22, f"{count} opps", size=8.5, color=MUTED)
    add_box(slide, 0.55, 6.42, 12.15, 0.42, fill="14213D", line=GRID)
    add_text(slide, 0.78, 6.50, 11.65, 0.20, f"Coverage check: {len(open_opportunities)}/{len(open_opportunities)} qualified open opportunities are direct; distributor account is blank in both exports.", size=9.5, color=ORANGE, bold=True)

    # 9 — Key account plans
    slide = new_slide(prs, "Q4 2026 / 2027 account plans: convert now, then expand adjacencies", "05 / Accounts", 9)
    account_plans = [
        ["Outdu", "AI response + tracking", "Close 21 Oct architecture, BOM and evaluation gate", "Replicate video/vision pipeline across tracking + analytics"],
        ["Honeywell", "Radar + industrial/aerospace control", "Qualify 3.00M radar; close two Design-stage evidence packs", "Create radar/thermal-control solution campaign"],
        ["Juniper", "PCIe FPGA + PQC interface", "Move CFPGA through validation; qualify PQC from 0%", "Position Agilex 3 for security/control-plane refresh"],
        ["GE HealthCare", "CT, anesthesia, power control", "Recover 2026 SP-PDU gate; define CT/anesthesia evaluations", "Medical imaging + controller platform expansion"],
        ["Emerson", "RX3i controller", "Freeze dual-device architecture before 16 Dec", "Industrial automation controller modernization"],
        ["Philips", "CT + MRI UI", "Secure CT design evidence; define MRI UI success criteria", "Imaging workflow and low-power control adjacency"],
        ["Ciena", "NID + access/aggregation routers", "Prioritize five open platforms and close architecture gaps", "Secure-control and network-interface platform expansion"],
    ]
    add_table(
        slide,
        0.55,
        1.35,
        12.15,
        4.95,
        ["Account", "Current plays", "Q4 2026 proposed focus", "2027 proposed expansion"],
        account_plans,
        [1.18, 2.28, 4.10, 4.59],
        font_size=9.0,
        header_size=8.5,
    )
    add_text(slide, 0.65, 6.43, 11.9, 0.22, "These are proposed actions inferred from stage/date/product data; validate customer commitments and support history in the review.", size=9.5, color=MUTED)

    # 10 — Account ownership
    slide = new_slide(prs, "Account ownership: focus list plus allocated-customer coverage", "06 / Account coverage", 10)
    add_box(slide, 0.55, 1.35, 3.70, 4.95, fill="14213D", line=TEAL)
    add_text(slide, 0.82, 1.66, 3.15, 0.24, "PRIMARY ACCOUNTS", size=9, color=TEAL, bold=True)
    for index, account in enumerate(PRIMARY_ACCOUNTS):
        y = 2.06 + index * 0.43
        add_text(slide, 0.84, y, 0.22, 0.22, "●", size=8, color=TEAL)
        add_text(slide, 1.13, y - 0.01, 2.75, 0.28, account, size=11.5, bold=True)
    add_text(slide, 0.82, 4.32, 3.15, 0.24, "OTHER STRATEGIC / SUPPORTED", size=9, color=ORANGE, bold=True)
    for index, account in enumerate(OTHER_STRATEGIC_ACCOUNTS):
        y = 4.70 + index * 0.29
        add_text(slide, 0.84, y, 0.20, 0.20, "•", size=8, color=ORANGE)
        add_text(slide, 1.10, y - 0.01, 2.82, 0.24, account, size=9.5, bold=True)
    roster_rows = [
        [customer, location, segment, allocation_date, dfae]
        for customer, location, segment, allocation_date, dfae in ALLOCATION_ROSTER
    ]
    add_table(
        slide,
        4.52,
        1.35,
        8.18,
        4.95,
        ["Customer", "Location", "Segment", "Allocation", "DFAE"],
        roster_rows,
        [1.20, 1.48, 0.90, 1.42, 3.18],
        font_size=7.2,
        header_size=7.4,
    )
    add_text(slide, 0.65, 6.45, 11.9, 0.22, "Allocation image: all rows list Sahil as New FAE. “Outdo” normalized to Outdu; unknown DFAE shown as TBD.", size=8.8, color=MUTED)

    # 11 — Support execution
    slide = new_slide(prs, "External support execution: 16 substantive issue threads in 75 days", "06 / Account coverage", 11)
    category_counts = Counter(issue[3] for issue in SUPPORT_ISSUES)
    resolved_count = sum(issue[5] == "Resolved" for issue in SUPPORT_ISSUES)
    add_stat_card(slide, 0.55, 1.35, 2.82, "Issue threads", str(len(SUPPORT_ISSUES)), "15 May–28 Jul 2026", TEAL)
    add_stat_card(slide, 3.55, 1.35, 2.82, "Coverage window", "75 days", "external technical email", BLUE)
    add_stat_card(slide, 6.55, 1.35, 2.82, "Technical themes", str(len(category_counts)), "from tools to board bring-up", ORANGE)
    add_stat_card(slide, 9.55, 1.35, 2.82, "Explicitly resolved", str(resolved_count), "closure proven in summaries", GREEN)
    category_colors = {
        "Board / Schematic": RED,
        "Networking": TEAL,
        "PCIe / Configuration": BLUE,
        "Tools / AI": ORANGE,
        "Device / I/O": GREEN,
    }
    max_count = max(category_counts.values())
    for index, (category, count) in enumerate(category_counts.most_common()):
        y = 2.95 + index * 0.60
        add_text(slide, 0.65, y + 0.02, 2.15, 0.22, category, size=10, bold=True)
        add_box(slide, 2.95, y, 6.65, 0.33, fill="152446", line="152446", radius=False)
        bar_width = 6.65 * count / max_count
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.95), Inches(y), Inches(bar_width), Inches(0.33))
        bar.fill.solid()
        bar.fill.fore_color.rgb = rgb(category_colors[category])
        bar.line.fill.background()
        add_text(slide, 9.78, y - 0.01, 0.45, 0.22, str(count), size=11, bold=True, align=PP_ALIGN.RIGHT)
    add_box(slide, 10.45, 2.90, 2.25, 3.03, fill="14213D", line=ORANGE)
    add_rich_text(
        slide,
        10.66,
        3.16,
        1.83,
        2.48,
        [
            ("QUALITY RULE", ORANGE, 9, True),
            ("Only one thread explicitly records customer-confirmed resolution.", WHITE, 11, True),
            ("Others show guidance, workaround, escalation or pending validation—not assumed closure.", PALE, 10, False),
        ],
    )
    add_citation(slide, "Source: user-provided distinct external email-thread analysis, 15 May–28 Jul 2026. Internal-only/admin threads excluded.")

    # 12 — Support-to-growth conversion
    slide = new_slide(prs, "Convert support into growth: primary-account actions", "06 / Account coverage", 12)
    def group_value(group: str) -> float:
        return sum(
            row["Peak Value"]
            for row in open_opportunities
            if account_in_group(group, str(row["Account Name"]))
        )

    account_actions = {
        "Ciena": ("5 qualified network platforms", "Joint OTN/router architecture and migration review"),
        "Juniper / HPE": ("3 PCIe/flash support threads", "Convert support into PCIe + secure-update design workshop"),
        "Philips": ("CT + MRI pipeline; no issue thread", "Secure CT evidence and MRI UI evaluation criteria"),
        "GE HealthCare": ("3 medical-control plays; no issue thread", "Run imaging/controller platform review before issues arise"),
        "Outdu": ("2 AI/vision plays; no issue thread", "Close AI architecture gate + run vision benchmark"),
    }
    primary_rows = []
    for account in PRIMARY_ACCOUNTS:
        support_threads = sum(issue_in_group(account, issue[1]) for issue in SUPPORT_ISSUES)
        signal, action = account_actions[account]
        primary_rows.append([account, fmt_value(group_value(account)), str(support_threads), signal, action])
    add_table(
        slide,
        0.55,
        1.35,
        12.15,
        4.55,
        ["Primary account", "Qualified open", "Threads", "Current signal", "Next growth action"],
        primary_rows,
        [1.45, 1.05, 0.65, 3.60, 5.40],
        font_size=9.2,
        header_size=8.4,
    )
    add_box(slide, 0.55, 6.12, 12.15, 0.55, fill="2B1D2F", line=ORANGE)
    add_text(slide, 0.78, 6.24, 11.62, 0.28, "Critical gap: 4 of 5 primary accounts had no substantive issue thread in this window. Move from reactive support to scheduled architecture creation.", size=10.2, color=ORANGE, bold=True)

    # 13 — Demand creation
    slide = new_slide(prs, "FY2027 demand creation: three repeatable motions", "07 / 2027 creation", 13)
    initiatives = [
        (
            "ROBOTICS + CONTROL",
            ORANGE,
            "Target",
            "Industrial OEMs; Emerson + Honeywell adjacencies",
            "Offer",
            "Deterministic control, functional partitioning, secure connectivity",
            "Create",
            "Controller architecture workshop + 3 qualified use cases",
        ),
        (
            "VIDEO + VISION",
            TEAL,
            "Target",
            "Outdu, medical imaging and smart sensing teams",
            "Offer",
            "Edge AI/video pipeline on Agilex 5 D / Agilex 3",
            "Create",
            "Vision benchmark day + reference design + 3 evaluations",
        ),
        (
            "INDUSTRIAL PLATFORM",
            BLUE,
            "Target",
            "Automation, medical controls and network infrastructure",
            "Offer",
            "PCIe/Ethernet modernization and legacy FPGA migration",
            "Create",
            "Portfolio migration clinic + joint Sales/FAE account map",
        ),
    ]
    for index, initiative in enumerate(initiatives):
        title, accent_color, label1, value1, label2, value2, label3, value3 = initiative
        x = 0.55 + index * 4.10
        add_box(slide, x, 1.42, 3.82, 4.75, fill="14213D", line=accent_color)
        add_text(slide, x + 0.22, 1.72, 3.38, 0.36, title, size=15, color=accent_color, bold=True)
        add_text(slide, x + 0.22, 2.35, 3.38, 0.22, label1.upper(), size=8, color=MUTED, bold=True)
        add_text(slide, x + 0.22, 2.62, 3.38, 0.70, value1, size=12, bold=True)
        add_text(slide, x + 0.22, 3.46, 3.38, 0.22, label2.upper(), size=8, color=MUTED, bold=True)
        add_text(slide, x + 0.22, 3.73, 3.38, 0.70, value2, size=12, bold=True)
        add_text(slide, x + 0.22, 4.60, 3.38, 0.22, label3.upper(), size=8, color=MUTED, bold=True)
        add_text(slide, x + 0.22, 4.87, 3.38, 0.86, value3, size=12, bold=True)
    add_text(slide, 0.65, 6.40, 11.9, 0.28, "Review decision: choose numeric 2027 creation targets after the 2027 sales objective is confirmed.", size=10.5, color=ORANGE, bold=True)

    # 14 — Critical review verdict
    slide = new_slide(prs, "Critical-review verdict: strong analysis, incomplete execution system", "08 / Readiness", 14)
    add_box(slide, 0.55, 1.35, 3.10, 4.98, fill="2B1D2F", line=RED)
    add_text(slide, 0.84, 1.70, 2.52, 0.24, "CURRENT VERDICT", size=9, color=RED, bold=True)
    add_text(slide, 0.84, 2.18, 2.52, 0.90, "NOT YET\nREADY TO SCALE", size=25, color=WHITE, bold=True)
    add_text(slide, 0.84, 3.38, 2.52, 1.72, "The market case and opportunity analysis are strong. The team, demos, distributor ownership and conversion system are not yet committed.", size=13, color=PALE, bold=True)
    add_text(slide, 0.84, 5.52, 2.52, 0.40, "Fix the operating model before launching campaigns.", size=11, color=ORANGE, bold=True)
    readiness_rows = [
        ["Pipeline facts", "GREEN", f"{len(open_opportunities)} qualified open + {len(outside_open_opportunities)} outside-export plays"],
        ["Market + competition", "GREEN", "Three sourced segments and an honest AMD battlecard"],
        ["Support execution", "GREEN", f"{len(SUPPORT_ISSUES)} external issue threads captured with outcomes"],
        ["Solution portfolio", "AMBER", "Solutions exist; access, kit, maturity and demo owners need verification"],
        ["Distributor team", "RED", "No Arrow or Macnica specialist is named"],
        ["Channel pipeline", "RED", f"0/{len(open_opportunities)} qualified open opportunities linked to a distributor"],
        ["Execution scorecard", "AMBER", "Proposed metrics exist; leadership and distis have not committed"],
    ]
    readiness_fills = ["14362E", "14362E", "14362E", "3A2F16", "3B1D28", "3B1D28", "3A2F16"]
    add_table(
        slide,
        3.92,
        1.35,
        8.78,
        4.98,
        ["Readiness gate", "Status", "Evidence / gap"],
        readiness_rows,
        [1.82, 0.90, 6.06],
        font_size=9.2,
        header_size=8.5,
        row_fills=readiness_fills,
    )
    add_text(slide, 3.98, 6.47, 8.60, 0.22, "Review decision: approve the joint Altera–Arrow–Macnica tiger team, owners, readiness dates and Q4 scorecard.", size=9.8, color=ORANGE, bold=True)

    # 15 — Altera solution stack
    slide = new_slide(prs, "Altera solution arsenal: lead with demonstrable systems, not device slides", "09 / Solution stack", 15)
    solution_columns = [
        (
            "ROBOTICS + CONTROL",
            ORANGE,
            [
                "ROS Consolidated Robot Controller",
                "Drive-on-Chip motor control",
                "Drive-on-Chip with PLC",
                "HPS TSN 3×2.5G SGMII",
                "Sensor Fusion Platform for AMRs",
                "Functional-safety flow / FSDP",
            ],
        ),
        (
            "CAMERA + EDGE AI",
            TEAL,
            [
                "Holoscan Sensor Bridge: MIPI→10GbE",
                "4Kp60 Multi-Sensor HDR Camera",
                "4Kp30 Multi-Sensor Camera + AI",
                "Smart Camera Demo Kit",
                "45+ Video & Vision Processing IPs",
                "FPGA AI Suite + MIPI CSI-2",
            ],
        ),
        (
            "PARTNER PLATFORMS",
            BLUE,
            [
                "Arrow Eagle Board",
                "Macnica Sulphur Agilex 5 kit",
                "Macnica MEP100 ST2110 SmartNIC",
                "Agilex 5 Modular / Premium kits",
                "Critical Link vision modules",
                "IntoPIX JPEG-XS ecosystem",
            ],
        ),
    ]
    for index, (title, accent_color, items) in enumerate(solution_columns):
        x = 0.55 + index * 4.10
        add_box(slide, x, 1.38, 3.82, 4.95, fill="14213D", line=accent_color)
        add_text(slide, x + 0.22, 1.70, 3.38, 0.28, title, size=12, color=accent_color, bold=True)
        for item_index, item in enumerate(items):
            y = 2.22 + item_index * 0.60
            add_text(slide, x + 0.22, y, 0.24, 0.24, "●", size=8, color=accent_color, bold=True)
            add_text(slide, x + 0.51, y - 0.01, 3.00, 0.45, item, size=10.5, bold=True)
        add_text(slide, x + 0.22, 5.93, 3.35, 0.20, "VERIFY: access • kit • owner • benchmark", size=8, color=MUTED, bold=True)
    add_citation(slide, "[R4], [R7], [R8], [R19], [R20]. Holoscan is NVIDIA technology integrated by Altera reference designs; solution access/maturity must be checked.")

    # 16 — Solution readiness plan
    slide = new_slide(prs, "Solution readiness: every offer needs an owner, running kit and proof", "09 / Solution stack", 16)
    solution_rows = [
        ["ROS robot controller", "Sahil + Robotics DFAE", "Agilex 5 SoC", "Sensor→actuator jitter + ROS 2 integration", "□"],
        ["Holoscan bridge", "Camera/AI DFAE + Macnica", "A5 Premium + MIPI + 10GbE", "4K stream into NVIDIA Holoscan", "□"],
        ["4Kp60 HDR camera", "Sahil + Macnica expert", "A5 Modular + camera", "ISP image quality, latency and power", "□"],
        ["4Kp30 camera + AI", "AI DFAE + Arrow expert", "A5 kit + customer model", "Ingest→inference→display benchmark", "□"],
        ["Drive-on-Chip", "Control DFAE + Arrow", "A5 Modular", "Motor-loop latency and integration", "□"],
        ["3×2.5G TSN", "Industrial DFAE + Sahil", "Agilex 5 SoC", "Multi-node sync, jitter and topology", "□"],
        ["Sensor fusion / AMR", "Robotics expert + Macnica", "A5 + sensor set", "Time-aligned multi-sensor pipeline", "□"],
        ["Smart camera", "Critical Link + joint team", "Partner kit", "Customer model + production BOM", "□"],
    ]
    add_table(
        slide,
        0.55,
        1.32,
        12.15,
        5.35,
        ["Solution", "Minimum owner", "Platform / kit", "Proof required before customer campaign", "Ready"],
        solution_rows,
        [1.62, 2.25, 2.18, 5.45, 0.65],
        font_size=8.3,
        header_size=8.0,
    )
    add_citation(slide, "Readiness means the team can run the demo, explain limitations, reproduce results and support an evaluation. A web page alone is not readiness.")

    # 17 — Joint team
    slide = new_slide(prs, "This is not a one-person job: form one market-penetration tiger team", "10 / Joint GTM", 17)
    team_columns = [
        (
            "ALTERA",
            TEAL,
            "Sahil — program + solution lead",
            [
                "Own segment strategy and architecture",
                "Coordinate Robotics, AI/Vision, TSN/FuSa DFAEs",
                "Define benchmark method and technical red lines",
                "Escalate product/IP/roadmap gaps",
            ],
            "NAMED",
        ),
        (
            "ARROW",
            ORANGE,
            "[Name required] — joint specialist",
            [
                "Robotics/AI/camera technical coverage",
                "Own Arrow account map and opportunity follow-up",
                "Maintain Eagle/A5 kits, samples and logistics",
                "Run workshops and CRM hygiene with Sales",
            ],
            "DUE 05 AUG",
        ),
        (
            "MACNICA",
            BLUE,
            "[Name required] — joint specialist",
            [
                "Robotics/AI/camera technical coverage",
                "Own Macnica account map and opportunity follow-up",
                "Bring Sulphur, MEP100 and video ecosystem",
                "Run evaluations and close integration gaps",
            ],
            "DUE 05 AUG",
        ),
    ]
    for index, (org, accent_color, lead, duties, status) in enumerate(team_columns):
        x = 0.55 + index * 4.10
        add_box(slide, x, 1.38, 3.82, 4.95, fill="14213D", line=accent_color)
        add_text(slide, x + 0.22, 1.67, 1.80, 0.28, org, size=14, color=accent_color, bold=True)
        add_text(slide, x + 2.16, 1.69, 1.37, 0.22, status, size=7.5, color=accent_color, bold=True, align=PP_ALIGN.RIGHT)
        add_text(slide, x + 0.22, 2.12, 3.38, 0.58, lead, size=13, bold=True)
        for duty_index, duty in enumerate(duties):
            y = 2.95 + duty_index * 0.64
            add_text(slide, x + 0.22, y, 0.22, 0.24, "●", size=8, color=accent_color)
            add_text(slide, x + 0.51, y - 0.01, 3.02, 0.48, duty, size=10.2, color=PALE, bold=True)
    add_text(slide, 0.65, 6.44, 11.9, 0.24, "Minimum formation rule: one named robotics/AI/camera-capable specialist from each organization, backed by Altera specialist DFAEs.", size=10, color=ORANGE, bold=True)

    # 18 — Joint penetration plan
    slide = new_slide(prs, "Joint market penetration: ready → target → prove → convert", "10 / Joint GTM", 18)
    phases = [
        ("BY 05 AUG", "FORM", RED, "Name Arrow + Macnica specialists; approve charter, segment roles and weekly cadence."),
        ("BY 14 AUG", "READY", ORANGE, "Run internal enablement; verify solution access; establish ≥1 working robotics and ≥1 camera/AI demo path."),
        ("BY 31 AUG", "TARGET", BLUE, "Build 30-account map: 15 Arrow + 15 Macnica; rank pains, installed competition and sponsor."),
        ("SEP–OCT", "ENGAGE", TEAL, "12 customer workshops; 8 benchmark/evaluation starts; joint Sales–FAE–disti follow-up in 48 hours."),
        ("Q4 CLOSE", "CONVERT", GREEN, "Proposed target: 6 qualified opportunities, 3 Develop/Design-stage plays and 2 DWINs."),
    ]
    for index, (when, name, accent_color, action) in enumerate(phases):
        y = 1.35 + index * 0.98
        add_box(slide, 0.55, y, 1.50, 0.74, fill="14213D", line=accent_color)
        add_text(slide, 0.74, y + 0.12, 1.12, 0.18, when, size=7.5, color=MUTED, bold=True)
        add_text(slide, 0.74, y + 0.37, 1.12, 0.22, name, size=11, color=accent_color, bold=True)
        add_box(slide, 2.25, y, 10.45, 0.74, fill=NAVY_2, line=GRID)
        add_text(slide, 2.55, y + 0.14, 9.85, 0.46, action, size=11.4, bold=True)
    add_box(slide, 0.55, 6.35, 12.15, 0.38, fill="2B1D2F", line=ORANGE)
    add_text(slide, 0.78, 6.42, 11.62, 0.20, "These are proposed critical-review commitments—not achieved results. Leadership, Arrow and Macnica must accept or reset them in the meeting.", size=8.8, color=ORANGE, bold=True)

    # 19 — Market signals
    slide = new_slide(prs, "Market review: growth is real, but proof-of-value wins budgets", "11 / Market review", 19)
    market_signals = [
        (
            "ROBOTICS",
            ORANGE,
            "9,100",
            "robots installed in India in 2024",
            "+7% YoY • India ranked #6 • automotive = 45% of installs",
            "Commercial implication",
            "Prioritize automotive suppliers, machine builders and control platforms; sell deterministic integration, not an FPGA.",
        ),
        (
            "VIDEO + VISION",
            TEAL,
            "$5.7B",
            "2025 global machine-vision forecast",
            "+1.5% after a 3.9% decline in 2024 • price pressure remains",
            "Commercial implication",
            "Lead with workload benchmarks, integration time and sensor flexibility; generic camera growth is not enough.",
        ),
        (
            "INDUSTRIAL",
            BLUE,
            "94%",
            "of APAC manufacturers investing/planning AI",
            "quality control 47% • process optimization 43% • cyber is critical",
            "Commercial implication",
            "Package edge AI with control, networking, safety and security around a measurable factory KPI.",
        ),
    ]
    for index, (name, accent_color, metric, metric_label, evidence, implication_label, implication) in enumerate(market_signals):
        x = 0.55 + index * 4.10
        add_box(slide, x, 1.38, 3.82, 4.95, fill="14213D", line=accent_color)
        add_text(slide, x + 0.22, 1.68, 3.38, 0.27, name, size=11, color=accent_color, bold=True)
        add_text(slide, x + 0.22, 2.13, 3.38, 0.54, metric, size=29, color=WHITE, bold=True)
        add_text(slide, x + 0.22, 2.72, 3.38, 0.46, metric_label, size=11, color=PALE, bold=True)
        add_text(slide, x + 0.22, 3.40, 3.38, 0.68, evidence, size=10.5, color=MUTED)
        add_text(slide, x + 0.22, 4.32, 3.38, 0.22, implication_label.upper(), size=8, color=accent_color, bold=True)
        add_text(slide, x + 0.22, 4.66, 3.38, 1.15, implication, size=11.5, bold=True)
    add_citation(slide, "[R1] IFR World Robotics 2025  •  [R3] Interact Analysis, Jun 2025  •  [R2] Rockwell APAC Smart Manufacturing 2025")

    # 20 — Robotics market
    slide = new_slide(prs, "Robotics: India is scaling; own the deterministic control layer", "11 / Market review", 20)
    add_box(slide, 0.55, 1.38, 4.00, 4.95, fill="14213D", line=ORANGE)
    add_text(slide, 0.80, 1.70, 3.50, 0.25, "MARKET EVIDENCE", size=9, color=ORANGE, bold=True)
    add_rich_text(
        slide,
        0.76,
        2.12,
        3.55,
        3.62,
        [
            ("542K global robot installs in 2024", WHITE, 18, True),
            ("4.664M operational stock; +9% YoY", PALE, 12, False),
            ("India: 9.1K installs, +7%, #6 globally", WHITE, 15, True),
            ("45% of India installs were automotive", PALE, 12, False),
            ("IFR forecasts 575K global installs in 2025 and >700K by 2028.", WHITE, 12, False),
        ],
    )
    add_box(slide, 4.78, 1.38, 3.78, 4.95, fill="14213D", line=TEAL)
    add_text(slide, 5.03, 1.70, 3.28, 0.25, "COMPETITIVE BATTLEFIELD", size=9, color=TEAL, bold=True)
    add_rich_text(
        slide,
        4.99,
        2.12,
        3.30,
        3.70,
        [
            ("AMD Xilinx strength", ORANGE, 10, True),
            ("Kria KR260: production SOM path, native ROS 2 and Kria Robotics Stack.", WHITE, 12, False),
            ("Altera proof", TEAL, 10, True),
            ("ROS 2 controller, Drive-on-Chip, safety, sensor fusion and 3×2.5G TSN reference designs.", WHITE, 12, False),
            ("Win condition", BLUE, 10, True),
            ("Demonstrate lower end-to-end jitter and fewer devices for the customer's exact control loop.", WHITE, 12, True),
        ],
    )
    add_box(slide, 8.78, 1.38, 3.92, 4.95, fill="14213D", line=BLUE)
    add_text(slide, 9.03, 1.70, 3.42, 0.25, "PROPOSED 2027 MOTION", size=9, color=BLUE, bold=True)
    add_rich_text(
        slide,
        8.99,
        2.12,
        3.45,
        3.78,
        [
            ("Target", BLUE, 10, True),
            ("Robot OEMs, AMR makers, automotive suppliers and control-system integrators.", WHITE, 12, False),
            ("Offer", BLUE, 10, True),
            ("Sense→think→act workshop using Agilex 5 SoC, TSN and functional-safety architecture.", WHITE, 12, False),
            ("Proposed scorecard", BLUE, 10, True),
            ("12 named accounts • 6 workshops • 3 evaluations • 2 qualified opportunities • 1 DWIN.", WHITE, 12, True),
        ],
    )
    add_citation(slide, "[R1] IFR World Robotics 2025  •  [R4] Altera Robotics Stack  •  [R5] Agilex 5  •  [R6] AMD Kria KR260")

    # 21 — Video and vision market
    slide = new_slide(prs, "Video + vision: recovery is selective; benchmark the whole pipeline", "11 / Market review", 21)
    add_box(slide, 0.55, 1.38, 4.00, 4.95, fill="14213D", line=TEAL)
    add_text(slide, 0.80, 1.70, 3.50, 0.25, "MARKET EVIDENCE", size=9, color=TEAL, bold=True)
    add_rich_text(
        slide,
        0.76,
        2.12,
        3.55,
        3.65,
        [
            ("2024 global market: −3.9%", WHITE, 18, True),
            ("Inventory correction and manufacturing softness", PALE, 12, False),
            ("2025 forecast: +1.5% to $5.7B", WHITE, 16, True),
            ("2028 forecast: $7B", PALE, 12, False),
            ("Area-scan cameras fell 7.8%; APAC vendor price pressure makes undifferentiated hardware vulnerable.", WHITE, 12, False),
        ],
    )
    add_box(slide, 4.78, 1.38, 3.78, 4.95, fill="14213D", line=ORANGE)
    add_text(slide, 5.03, 1.70, 3.28, 0.25, "COMPETITIVE BATTLEFIELD", size=9, color=ORANGE, bold=True)
    add_rich_text(
        slide,
        4.99,
        2.12,
        3.30,
        3.72,
        [
            ("AMD Xilinx strength", ORANGE, 10, True),
            ("KV260 quick start, multi-camera interfaces, Vitis Vision/AI and Versal AI Engine scale.", WHITE, 12, False),
            ("Altera proof", TEAL, 10, True),
            ("45+ VVP IP cores, >600 MHz claimed Fmax, 8K60+ optimization and FPGA AI Suite.", WHITE, 12, False),
            ("Win condition", BLUE, 10, True),
            ("Customer-model benchmark: ingest + ISP + inference + output latency, power and engineering weeks.", WHITE, 12, True),
        ],
    )
    add_box(slide, 8.78, 1.38, 3.92, 4.95, fill="14213D", line=BLUE)
    add_text(slide, 9.03, 1.70, 3.42, 0.25, "PROPOSED 2027 MOTION", size=9, color=BLUE, bold=True)
    add_rich_text(
        slide,
        8.99,
        2.12,
        3.45,
        3.78,
        [
            ("Named beachheads", BLUE, 10, True),
            ("Outdu tracking/analytics; GE + Philips medical imaging; industrial inspection.", WHITE, 12, False),
            ("Offer", BLUE, 10, True),
            ("Bring-your-model vision benchmark day with two sensor interfaces and a production BOM.", WHITE, 12, False),
            ("Proposed scorecard", BLUE, 10, True),
            ("8 accounts • 4 benchmarks • 3 evaluations • 2 qualified opportunities • 1 DWIN.", WHITE, 12, True),
        ],
    )
    add_citation(slide, "[R3] Interact Analysis  •  [R7] Altera VVP Suite  •  [R8] FPGA AI Suite  •  [R9] AMD KV260  •  [R10] Versal AI Edge")

    # 22 — Industrial market
    slide = new_slide(prs, "Industrial: AI spending converges with control, safety and cyber", "11 / Market review", 22)
    add_box(slide, 0.55, 1.38, 4.00, 4.95, fill="14213D", line=BLUE)
    add_text(slide, 0.80, 1.70, 3.50, 0.25, "MARKET EVIDENCE — APAC", size=9, color=BLUE, bold=True)
    add_rich_text(
        slide,
        0.76,
        2.12,
        3.55,
        3.65,
        [
            ("94% investing / planning AI", WHITE, 18, True),
            ("47%: quality control is top AI use case", PALE, 12, False),
            ("43%: process optimization", WHITE, 16, True),
            ("95%: cyber standards important", PALE, 12, False),
            ("The buying unit spans OT control, data, safety and security—not a single component owner.", WHITE, 12, False),
        ],
    )
    add_box(slide, 4.78, 1.38, 3.78, 4.95, fill="14213D", line=ORANGE)
    add_text(slide, 5.03, 1.70, 3.28, 0.25, "COMPETITIVE BATTLEFIELD", size=9, color=ORANGE, bold=True)
    add_rich_text(
        slide,
        4.99,
        2.12,
        3.30,
        3.72,
        [
            ("Shared table stakes", ORANGE, 10, True),
            ("AMD and Altera both offer TSN, industrial connectivity and certified safety flows.", WHITE, 12, False),
            ("Altera proof point", TEAL, 10, True),
            ("Agilex 5 SoC integrates three hardened 2.5G TSN MACs; Drive-on-Chip and FSDP support.", WHITE, 12, False),
            ("Win condition", BLUE, 10, True),
            ("Quantify BOM consolidation, deterministic cycle time and certification work saved.", WHITE, 12, True),
        ],
    )
    add_box(slide, 8.78, 1.38, 3.92, 4.95, fill="14213D", line=TEAL)
    add_text(slide, 9.03, 1.70, 3.42, 0.25, "PROPOSED 2027 MOTION", size=9, color=TEAL, bold=True)
    add_rich_text(
        slide,
        8.99,
        2.12,
        3.45,
        3.78,
        [
            ("Named beachheads", TEAL, 10, True),
            ("Emerson controller; Honeywell automation; machine builders via Arrow/Macnica.", WHITE, 12, False),
            ("Offer", TEAL, 10, True),
            ("Industrial platform clinic: TSN + PLC/control + safety + secure lifecycle.", WHITE, 12, False),
            ("Proposed scorecard", TEAL, 10, True),
            ("15 accounts • 6 clinics • 4 evaluations • 3 qualified opportunities • 1 DWIN.", WHITE, 12, True),
        ],
    )
    add_citation(slide, "[R2] Rockwell APAC 2025  •  [R11–R15] Altera/AMD industrial networking and functional-safety sources")

    # 23 — Portfolio map
    slide = new_slide(prs, "Competitive portfolio map: no single one-for-one device comparison", "12 / Competition", 23)
    portfolio_rows = [
        [
            "Cost / I/O edge",
            "Agilex 3 C: FPGA + optional dual A55 SoC, AI tensor DSP, MIPI, 12.5G",
            "Spartan UltraScale+: high I/O, 16.3G, PCIe Gen4, hardened LPDDR; no processor",
            "Compare package, I/O mix, BOM, power and required processor—not density labels.",
        ],
        [
            "Midrange embedded",
            "Agilex 5 E: A76/A55 SoC options, 3×2.5G TSN, AI tensor blocks",
            "Zynq UltraScale+ MPSoC / Kria K26: mature SOM and software ecosystem",
            "Altera: consolidation + hard TSN. AMD: fast SOM/ROS path. Prove development effort.",
        ],
        [
            "AI / vision scale",
            "Agilex 5 D/E + FPGA AI Suite + VVP; custom streaming data path",
            "Versal AI Edge + AI Engines; Kria KV260 for fast vision evaluation",
            "TOPS are not directly comparable. Benchmark the customer's model and full pipeline.",
        ],
        [
            "High-end adaptable",
            "Agilex 7/9 families for high bandwidth, RF and complex acceleration",
            "Versal Premium / AI Core / RF families",
            "Treat as architecture-led pursuits; involve product-line DFAEs before device claims.",
        ],
    ]
    add_table(
        slide,
        0.55,
        1.35,
        12.15,
        4.95,
        ["Battlefield", "Altera position", "AMD Xilinx position", "FAE comparison rule"],
        portfolio_rows,
        [1.35, 3.30, 3.30, 4.20],
        font_size=8.6,
        header_size=8.4,
    )
    add_text(slide, 0.65, 6.43, 11.9, 0.22, "Never compare LE vs logic-cell counts or headline TOPS as equivalent metrics; normalize the workload, precision, sparsity, clocks, power and tool version.", size=9.5, color=ORANGE, bold=True)
    add_citation(slide, "[R5], [R8], [R10], [R16], [R17] — official vendor product pages; claims remain vendor-stated until customer benchmark")

    # 24 — Competitive scorecard
    slide = new_slide(prs, "Altera vs AMD Xilinx: evidence-based battlecard", "12 / Competition", 24)
    score_rows = [
        ["Robotics SW", "ROS 2 controller + Drive-on-Chip reference stack", "KR260 + native ROS 2 + KRS + production K26 SOM", "AMD advantage for software-first/SOM buyers; counter with deterministic consolidation proof."],
        ["Vision IP", "45+ VVP cores; FPGA AI Suite; custom streaming fabric", "Vitis Vision/AI; KV260; Versal AI Engines", "Both credible. Win on measured end-to-end latency, power and engineering effort."],
        ["AI scale", "Agilex 3/5/7 portfolio; vendor claims up to 152.6 INT8 TOPS on A5 D", "Versal AI Edge vendor table: 5–202 dense INT8 TOPS", "Do not rank by TOPS alone; architectures, devices and assumptions differ."],
        ["Industrial TSN", "Agilex 5 SoC: three hardened 2.5G TSN MACs", "100M/1G TSN LogiCORE; broader industrial protocol ecosystem", "Lead with integration/BOM if hard TSN fits; validate exact protocol and topology."],
        ["Functional safety", "TÜV-reviewed FSDP/methodology; stated suitability up to SIL3", "TÜV SÜD-certified flows; IEC 61508 / ISO 13849 support", "Table stakes for both. Compare exact device/tool/IP certificate scope."],
        ["Cost / I/O", "Agilex 3: SoC option, MIPI, AI tensor DSP, 12.5G", "Spartan UltraScale+: up to 572 I/O, 16.3G, PCIe Gen4", "Use customer pinout + board-cost study. Neither wins every configuration."],
    ]
    add_table(
        slide,
        0.55,
        1.32,
        12.15,
        5.20,
        ["Criterion", "Altera evidence", "AMD Xilinx evidence", "Implication / honest position"],
        score_rows,
        [1.10, 3.08, 3.08, 4.89],
        font_size=7.6,
        header_size=7.8,
    )
    add_citation(slide, "[R4–R18] Official vendor sources. “Advantage” statements are sales assessments, not third-party benchmark conclusions.")

    # 25 — Competitive win plan
    slide = new_slide(prs, "How to win against AMD Xilinx: replace claims with customer evidence", "12 / Competition", 25)
    win_steps = [
        ("1  DISCOVER", ORANGE, "Installed device/tool flow? Decision owner? Workload? Pain: latency, power, I/O, BOM, safety or schedule?"),
        ("2  BASELINE", RED, "Capture AMD reference: exact device/board, tool version, model precision, interfaces, clocks and measured power."),
        ("3  PROVE", TEAL, "Run the same workload on Altera; publish reproducible scripts, resource use, latency distribution, power and BOM."),
        ("4  DE-RISK", BLUE, "Close IP gaps, migration effort, safety scope, supply/lifecycle and support plan with named DFAEs."),
        ("5  COMMERCIALIZE", GREEN, "Convert evaluation success into architecture freeze, DWIN evidence and production forecast."),
    ]
    for index, (label, accent_color, action) in enumerate(win_steps):
        y = 1.38 + index * 0.97
        add_box(slide, 0.55, y, 2.02, 0.72, fill="14213D", line=accent_color)
        add_text(slide, 0.78, y + 0.20, 1.58, 0.25, label, size=11, color=accent_color, bold=True)
        add_box(slide, 2.75, y, 9.95, 0.72, fill=NAVY_2, line=GRID)
        add_text(slide, 3.02, y + 0.14, 9.42, 0.42, action, size=11.5, bold=True)
    add_box(slide, 0.55, 6.38, 12.15, 0.35, fill="2B1D2F", line=ORANGE)
    add_text(slide, 0.78, 6.43, 11.62, 0.20, "Red line: do not use vendor headline TOPS, power or performance claims as an apples-to-apples result without reproducing the workload.", size=9, color=ORANGE, bold=True)

    # 26 — Ecosystem
    slide = new_slide(prs, "DFAE + ecosystem engagement: attach expertise to stage exits", "13 / Team execution", 26)
    columns = [
        ("DISCOVERY", "Sales + FAE", "Sponsor, use case, value, budget and decision process", RED),
        ("ARCHITECTURE", "FAE + specialist DFAE", "Block diagram, device fit, IP/tools and risk register", ORANGE),
        ("VALIDATION", "Customer + FAE + DFAE", "Evaluation results, issue closure and design-freeze evidence", BLUE),
        ("PRODUCTION", "Sales + FAE + channel", "DW evidence, forecast, supply path and production handoff", GREEN),
    ]
    for index, (stage, owners, evidence, accent_color) in enumerate(columns):
        x = 0.55 + index * 3.08
        add_box(slide, x, 1.48, 2.82, 2.62, fill="14213D", line=accent_color)
        add_text(slide, x + 0.18, 1.76, 2.46, 0.25, stage, size=11, color=accent_color, bold=True)
        add_text(slide, x + 0.18, 2.18, 2.46, 0.24, owners, size=12.5, bold=True)
        add_text(slide, x + 0.18, 2.72, 2.46, 0.94, evidence, size=11.5, color=PALE)
    add_box(slide, 0.55, 4.42, 12.15, 1.65)
    add_text(slide, 0.82, 4.72, 2.2, 0.22, "WEEKLY OPERATING CADENCE", size=9, color=TEAL, bold=True)
    cadence = [
        ("MON", "Top-5 stage/gap inspection"),
        ("WED", "Customer/DFAE action review"),
        ("FRI", "Evidence + CRM hygiene"),
        ("MONTH-END", "Pipeline creation + conversion scorecard"),
    ]
    for index, (when, action) in enumerate(cadence):
        x = 3.00 + index * 2.35
        add_text(slide, x, 4.70, 1.95, 0.20, when, size=8, color=MUTED, bold=True)
        add_text(slide, x, 5.00, 2.05, 0.62, action, size=11, bold=True)

    # 27 — Commitments / gaps
    slide = new_slide(prs, "Close the review with owners, evidence and missing inputs", "14 / Commitments", 27)
    add_box(slide, 0.55, 1.38, 7.30, 4.95)
    add_text(slide, 0.82, 1.70, 6.75, 0.26, "PROPOSED COMMITMENTS", size=10, color=TEAL, bold=True)
    commitments = [
        "Name next-stage evidence and due date for every top-5 opportunity.",
        "Qualify Weather Radar and CPU Interface Card from 0% or remove value from the active forecast.",
        "Name one Arrow and one Macnica robotics/AI/camera specialist by 05 Aug; verify solution readiness by 14 Aug.",
        "Publish a 30-account Arrow + Macnica map and create the 15th strategic opportunity.",
        "Run three joint benchmark-led campaigns and inspect workshops, evaluations, qualified pipeline and DWINs weekly.",
    ]
    for index, commitment in enumerate(commitments):
        y = 2.18 + index * 0.76
        add_text(slide, 0.84, y, 0.35, 0.30, f"{index + 1}", size=15, color=ORANGE, bold=True)
        add_text(slide, 1.28, y, 6.10, 0.55, commitment, size=13, bold=True)
    add_box(slide, 8.10, 1.38, 4.60, 4.95, fill="14213D", line=RED)
    add_text(slide, 8.38, 1.70, 4.05, 0.26, "INPUTS TO COMPLETE BEFORE PRESENTING", size=10, color=RED, bold=True)
    missing = [
        "2026 target",
        "YTD actual / achieved DWINs",
        "Gap-closure value and dates",
        "Support activity history",
        "Customer market-share baseline",
        "Arrow + Macnica specialist names",
        "Named solution DFAEs / kit status",
    ]
    for index, item in enumerate(missing):
        y = 2.18 + index * 0.49
        add_text(slide, 8.40, y, 0.28, 0.22, "□", size=13, color=ORANGE, bold=True)
        add_text(slide, 8.78, y, 3.35, 0.28, item, size=11.5, bold=True)
    add_text(slide, 8.40, 5.76, 3.95, 0.34, "No invented numbers.\nFill, validate, commit.", size=13.5, color=TEAL, bold=True)

    # 28 — Data definitions
    slide = new_slide(prs, "Appendix: scope, definitions and data-quality notes", "Appendix", 28)
    notes = [
        ("Scope", f"Reporting scope attributes Sahil Patni + Kasturi Rangan to Sahil. Master: {len(opportunities)} opportunities; open exports: {len(open_opportunities)}."),
        ("Buckets", "“2026” and “2027” are the user-provided export labels. Membership is not inferred from Design Win Date; the 2027 export includes some 2026 DWIN dates."),
        ("Peak value", "Product-line Peak Value summed within each Opportunity ID. Source file does not encode currency."),
        ("Weighted value", "Consolidated peak value × source probability. Identify opportunities at 0% therefore contribute zero."),
        ("DWIN-dated", "Opportunity Design Win Date falls in the stated year. This is a scheduled date, not proof of an achieved design win."),
        ("Dates", "Source mixes Excel dates and text dates; both were normalized for analysis."),
        ("Coverage", f"All {len(open_opportunities)} uploaded open opportunities are direct and Distributor Account is blank. No channel-owned pipeline is evidenced."),
    ]
    for index, (label, note) in enumerate(notes):
        y = 1.36 + index * 0.72
        add_text(slide, 0.65, y, 1.65, 0.26, label.upper(), size=9, color=TEAL if index < 4 else ORANGE, bold=True)
        add_text(slide, 2.20, y, 10.15, 0.48, note, size=11.5, color=WHITE)
    add_text(slide, 0.65, 6.56, 11.9, 0.24, "Regenerate with: python3 qbr/generate_sahil_qbr.py", size=9, color=MUTED)

    # 29/30 — Reference appendix
    for slide_number, start in ((29, 0), (30, 10)):
        end = start + 10
        title = "References: market evidence and solution sources" if start == 0 else "References: competitive product and tool sources"
        slide = new_slide(prs, title, "Appendix / References", slide_number)
        for index, (reference_id, reference_title, url) in enumerate(REFERENCES[start:end]):
            y = 1.25 + index * 0.53
            add_text(slide, 0.62, y, 0.55, 0.22, reference_id, size=9, color=TEAL, bold=True)
            add_text(slide, 1.15, y, 4.40, 0.36, reference_title, size=7.8, color=WHITE, bold=True)
            add_text(slide, 5.68, y, 6.92, 0.36, url, size=6.8, color=PALE)
        add_text(slide, 0.65, 6.70, 11.9, 0.18, "Accessed 28 Jul 2026. Vendor specifications and performance figures are vendor-stated; validate against current datasheets and customer workloads.", size=7.5, color=ORANGE, bold=True)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


def build_summary_workbook(
    line_items: list[dict[str, Any]],
    opportunities: list[dict[str, Any]],
    open_buckets: dict[str, list[dict[str, Any]]],
    output: Path,
):
    workbook = Workbook()
    summary = workbook.active
    summary.title = "Executive Summary"

    total = sum(row["Peak Value"] for row in opportunities)
    open_opportunities = open_buckets["2026"] + open_buckets["2027"]
    open_ids = {row["Opportunity ID"] for row in open_opportunities}
    outside_open_opportunities = [
        row for row in opportunities if row["Opportunity ID"] not in open_ids
    ]
    identify_opportunities = [
        row
        for row in outside_open_opportunities
        if row["Opportunity Stage"] == "Identify"
    ]
    other_outside_opportunities = [
        row
        for row in outside_open_opportunities
        if row["Opportunity Stage"] != "Identify"
    ]
    ranked_opportunities = sorted(
        open_opportunities, key=lambda row: row["Peak Value"], reverse=True
    ) + sorted(outside_open_opportunities, key=lambda row: row["Peak Value"], reverse=True)
    qualified_total = sum(row["Peak Value"] for row in open_opportunities)
    weighted = sum(row["Weighted Value"] for row in open_opportunities)
    metrics = [
        ("Metric", "Value", "Definition"),
        ("Qualified open opportunities", len(open_opportunities), "Combined uploaded 2026 + 2027 open exports"),
        ("Qualified open peak value", qualified_total, "Currency absent from source"),
        ("Qualified open weighted value", weighted, "Peak Value × normalized Probability"),
        ("2026 open-export opportunities", len(open_buckets["2026"]), "User-provided planning bucket"),
        ("2026 open-export peak value", sum(row["Peak Value"] for row in open_buckets["2026"]), "User-provided planning bucket"),
        ("2027 open-export opportunities", len(open_buckets["2027"]), "User-provided planning bucket"),
        ("2027 open-export peak value", sum(row["Peak Value"] for row in open_buckets["2027"]), "User-provided planning bucket"),
        ("Identify outside open exports", len(identify_opportunities), "Master-report opportunities absent from both open exports"),
        ("Identify peak value", sum(row["Peak Value"] for row in identify_opportunities), "0% probability in master report"),
        ("Other outside open exports", len(other_outside_opportunities), "Non-Identify master opportunities absent from both exports"),
        ("Other outside peak value", sum(row["Peak Value"] for row in other_outside_opportunities), "Review reason for exclusion"),
        ("Total discovered opportunities", len(opportunities), "Qualified open + outside-export plays"),
        ("Total discovered peak value", total, "Qualified open + outside-export plays"),
        ("Reporting scope", "Sahil + Kasturi", "Only these two source Technical Owner values are included"),
        ("Product line items", len(line_items), "Technical Owner = Sahil Patni or Kasturi Rangan"),
        ("External support issue threads", len(SUPPORT_ISSUES), "Distinct substantive external email threads, 15 May–28 Jul 2026"),
        ("Explicitly resolved support threads", sum(issue[5] == "Resolved" for issue in SUPPORT_ISSUES), "Closure explicitly recorded in supplied summaries"),
        ("2026 annual target", "INPUT REQUIRED", "Not present in source"),
        ("YTD actual achievement", "INPUT REQUIRED", "Not present in source"),
    ]
    for row in metrics:
        summary.append(row)

    opportunity_sheet = workbook.create_sheet("Opportunities")
    headers = [
        "Rank",
        "Open Pipeline Bucket",
        "Source Technical Owner",
        "Opportunity ID",
        "Account",
        "Opportunity",
        "Stage",
        "Probability (%)",
        "Design Win Date",
        "Created Date",
        "Production Start Date",
        "Vertical",
        "Peak Value",
        "Weighted Value",
        "Product Count",
        "Products",
        "Suggested Next Gate",
    ]
    opportunity_sheet.append(headers)
    bucket_by_id = {
        row["Opportunity ID"]: bucket
        for bucket, bucket_rows in open_buckets.items()
        for row in bucket_rows
    }
    for rank, row in enumerate(ranked_opportunities, start=1):
        opportunity_sheet.append(
            [
                rank,
                bucket_by_id.get(
                    row["Opportunity ID"],
                    "Identify / outside open exports"
                    if row["Opportunity Stage"] == "Identify"
                    else "Outside uploaded 2026/2027",
                ),
                row["Technical Owner"],
                row["Opportunity ID"],
                row["Account Name"],
                row["Opportunity Name"],
                row["Opportunity Stage"],
                row["Probability (%)"],
                row["Design Win Date Parsed"],
                row["Created Date Parsed"],
                row["Production Start Date Parsed"],
                row["Vertical Market"],
                row["Peak Value"],
                row["Weighted Value"],
                len(row["Products"]),
                "\n".join(str(product) for product in row["Products"]),
                suggested_gate(row),
            ]
        )

    reconciliation_sheet = workbook.create_sheet("Pipeline Reconciliation")
    reconciliation_sheet.append(["Population", "Opportunities", "Peak Value", "Weighted Value", "Interpretation"])
    reconciliation_rows = [
        [
            "2026 open export",
            len(open_buckets["2026"]),
            sum(row["Peak Value"] for row in open_buckets["2026"]),
            sum(row["Weighted Value"] for row in open_buckets["2026"]),
            "User-provided planning bucket; not derived from DWIN year",
        ],
        [
            "2027 open export",
            len(open_buckets["2027"]),
            sum(row["Peak Value"] for row in open_buckets["2027"]),
            sum(row["Weighted Value"] for row in open_buckets["2027"]),
            "User-provided planning bucket; includes some 2026 DWIN dates",
        ],
        [
            "Qualified open combined",
            len(open_opportunities),
            qualified_total,
            weighted,
            "No overlapping opportunity IDs across uploaded exports",
        ],
        [
            "Identify outside open exports",
            len(identify_opportunities),
            sum(row["Peak Value"] for row in identify_opportunities),
            sum(row["Weighted Value"] for row in identify_opportunities),
            "Weather Radar and CPU Interface Card; 0% probability",
        ],
        [
            "Other outside open exports",
            len(other_outside_opportunities),
            sum(row["Peak Value"] for row in other_outside_opportunities),
            sum(row["Weighted Value"] for row in other_outside_opportunities),
            "Data Recorder (Kasturi): Define / 25%; reason for exclusion requires review",
        ],
        [
            "Total discovered portfolio",
            len(opportunities),
            total,
            sum(row["Weighted Value"] for row in opportunities),
            "Master-report universe",
        ],
    ]
    for row in reconciliation_rows:
        reconciliation_sheet.append(row)

    open_headers = [
        "Rank",
        "Source Technical Owner",
        "Opportunity ID",
        "Account",
        "Opportunity",
        "Stage",
        "Probability (%)",
        "Design Win Date",
        "Peak Value",
        "Weighted Value",
    ]
    for bucket, bucket_rows in open_buckets.items():
        bucket_sheet = workbook.create_sheet(f"Open Pipeline {bucket}")
        bucket_sheet.append(open_headers)
        for rank, row in enumerate(bucket_rows, start=1):
            bucket_sheet.append(
                [
                    rank,
                    row["Technical Owner"],
                    row["Opportunity ID"],
                    row["Account Name"],
                    row["Opportunity Name"],
                    row["Opportunity Stage"],
                    row["Probability (%)"],
                    row["Design Win Date Parsed"],
                    row["Peak Value"],
                    row["Weighted Value"],
                ]
            )

    account_sheet = workbook.create_sheet("Account Coverage")
    account_sheet.append(
        [
            "Tier",
            "Account",
            "Qualified Open Value",
            "Outside-Export Value",
            "Support Threads",
            "Coverage Note",
        ]
    )
    account_notes = {
        "Ciena": "Primary; five qualified network platforms",
        "Juniper / HPE": "Primary; PCIe/flash support plus qualified and Identify pipeline",
        "Philips": "Primary; CT and MRI pipeline",
        "GE HealthCare": "Primary; CT, anesthesia and power-control pipeline",
        "Outdu": "Primary; AI response and tracking pipeline",
        "Siemens": "Other strategic; MAX 10 GPIO support",
        "Boeing": "Other strategic; Data Recorder sits outside uploaded open exports",
        "Honeywell Aerospace": "Other strategic; I/O support and Identify-stage radar",
        "Schneider": "Other strategic; no supplied pipeline/support thread in period",
        "Emerson": "Other strategic; RX3i pipeline",
    }
    for tier, accounts in (
        ("Primary", PRIMARY_ACCOUNTS),
        ("Other strategic / supported", OTHER_STRATEGIC_ACCOUNTS),
    ):
        for account in accounts:
            account_sheet.append(
                [
                    tier,
                    account,
                    sum(
                        row["Peak Value"]
                        for row in open_opportunities
                        if account_in_group(account, str(row["Account Name"]))
                    ),
                    sum(
                        row["Peak Value"]
                        for row in outside_open_opportunities
                        if account_in_group(account, str(row["Account Name"]))
                    ),
                    sum(issue_in_group(account, issue[1]) for issue in SUPPORT_ISSUES),
                    account_notes[account],
                ]
            )

    allocation_sheet = workbook.create_sheet("Allocation Roster")
    allocation_sheet.append(
        ["Customer", "Location", "Segment", "Allocation Date", "DFAE", "New FAE"]
    )
    for customer, location, segment, allocation_date, dfae in ALLOCATION_ROSTER:
        allocation_sheet.append(
            [customer, location, segment, allocation_date, dfae, "Sahil"]
        )

    support_sheet = workbook.create_sheet("External Support Issues")
    support_sheet.append(
        ["#", "Issue", "Customer / Partner", "Latest Activity", "Theme", "Technical Outcome / Next Action", "Status"]
    )
    for index, (issue, customer, latest, category, outcome, status) in enumerate(
        SUPPORT_ISSUES, start=1
    ):
        support_sheet.append(
            [index, issue, customer, latest, category, outcome, status]
        )

    source_sheet = workbook.create_sheet("Sahil Line Items")
    source_headers = list(line_items[0].keys())
    source_sheet.append(source_headers)
    for row in line_items:
        source_sheet.append([row[header] for header in source_headers])

    reference_sheet = workbook.create_sheet("Market References")
    reference_sheet.append(["ID", "Source", "URL", "Accessed", "Use"])
    reference_uses = {
        "R1": "Robotics market",
        "R2": "Industrial / APAC market",
        "R3": "Machine-vision market",
        "R4": "Altera robotics",
        "R5": "Altera Agilex 5",
        "R6": "AMD robotics",
        "R7": "Altera video / vision",
        "R8": "Altera AI software",
        "R9": "AMD vision",
        "R10": "AMD AI portfolio",
        "R11": "Altera industrial",
        "R12": "Altera TSN / HPS",
        "R13": "AMD industrial networking",
        "R14": "Altera functional safety",
        "R15": "AMD functional safety",
        "R16": "Altera cost-optimized portfolio",
        "R17": "AMD cost-optimized portfolio",
        "R18": "AMD development tools",
        "R19": "Altera video solution stack",
        "R20": "Altera Holoscan / sensor interfaces",
    }
    for reference_id, reference_title, url in REFERENCES:
        reference_sheet.append([reference_id, reference_title, url, date(2026, 7, 28), reference_uses[reference_id]])

    joint_sheet = workbook.create_sheet("Joint GTM Plan")
    joint_sheet.append(["Workstream", "Owner", "Commitment", "Due", "Proposed KPI", "Status"])
    joint_rows = [
        ["Team", "Altera / Sahil", "Lead program, architecture, benchmarks and DFAE coordination", date(2026, 8, 5), "Named and charter accepted", "NAMED"],
        ["Team", "Arrow", "Name robotics/AI/camera-capable joint specialist", date(2026, 8, 5), "≥1 named specialist", "OPEN"],
        ["Team", "Macnica", "Name robotics/AI/camera-capable joint specialist", date(2026, 8, 5), "≥1 named specialist", "OPEN"],
        ["Readiness", "Joint tiger team", "Verify solution access, kits, skills and benchmark playbooks", date(2026, 8, 14), "≥1 robotics + ≥1 camera/AI demo path", "OPEN"],
        ["Targeting", "Arrow + Macnica", "Create non-duplicated joint account map", date(2026, 8, 31), "30 accounts: 15 per distributor", "OPEN"],
        ["Engagement", "Joint tiger team", "Run segment workshops and launch evaluations", date(2026, 10, 31), "12 workshops; 8 evaluations", "OPEN"],
        ["Conversion", "Sales + FAE + distis", "Advance evidence-backed opportunities", date(2026, 12, 31), "6 qualified; 3 Develop/Design; 2 DWINs", "PROPOSED"],
    ]
    for row in joint_rows:
        joint_sheet.append(row)

    for sheet in workbook.worksheets:
        sheet.freeze_panes = "A2"
        sheet.auto_filter.ref = sheet.dimensions
        for cell in sheet[1]:
            cell.fill = PatternFill("solid", fgColor=BLUE)
            cell.font = Font(color="FFFFFF", bold=True)
            cell.alignment = Alignment(vertical="center")
        for row in sheet.iter_rows(min_row=2):
            for cell in row:
                cell.alignment = Alignment(vertical="top", wrap_text=True)
        for column_cells in sheet.columns:
            max_length = max(len(str(cell.value or "")) for cell in column_cells)
            sheet.column_dimensions[column_cells[0].column_letter].width = min(max(max_length + 2, 10), 44)
    for row in summary.iter_rows(min_row=2):
        if isinstance(row[1].value, (int, float)):
            row[1].number_format = '#,##0.00'
    for sheet in workbook.worksheets:
        header_by_column = {cell.column: cell.value for cell in sheet[1]}
        value_columns = {
            column
            for column, header in header_by_column.items()
            if header
            in {
                "Peak Value",
                "Weighted Value",
                "Qualified Open Value",
                "Outside-Export Value",
            }
        }
        for row in sheet.iter_rows(min_row=2):
            for cell in row:
                if cell.column in value_columns and isinstance(cell.value, (int, float)):
                    cell.number_format = '#,##0.00'
    opportunity_sheet.column_dimensions["P"].width = 52
    opportunity_sheet.column_dimensions["Q"].width = 48
    reconciliation_sheet.column_dimensions["E"].width = 56
    reference_sheet.column_dimensions["B"].width = 58
    reference_sheet.column_dimensions["C"].width = 90
    joint_sheet.column_dimensions["C"].width = 62
    joint_sheet.column_dimensions["E"].width = 42
    account_sheet.column_dimensions["F"].width = 58
    allocation_sheet.column_dimensions["B"].width = 18
    support_sheet.column_dimensions["B"].width = 54
    support_sheet.column_dimensions["C"].width = 25
    support_sheet.column_dimensions["F"].width = 76

    output.parent.mkdir(parents=True, exist_ok=True)
    workbook.save(output)


def build_notes(
    opportunities: list[dict[str, Any]],
    open_buckets: dict[str, list[dict[str, Any]]],
    output: Path,
):
    total = sum(row["Peak Value"] for row in opportunities)
    open_2026 = open_buckets["2026"]
    open_2027 = open_buckets["2027"]
    open_opportunities = open_2026 + open_2027
    open_ids = {row["Opportunity ID"] for row in open_opportunities}
    outside_open_opportunities = [
        row for row in opportunities if row["Opportunity ID"] not in open_ids
    ]
    identify_opportunities = [
        row
        for row in outside_open_opportunities
        if row["Opportunity Stage"] == "Identify"
    ]
    other_outside_opportunities = [
        row
        for row in outside_open_opportunities
        if row["Opportunity Stage"] != "Identify"
    ]
    qualified_total = sum(row["Peak Value"] for row in open_opportunities)
    weighted = sum(row["Weighted Value"] for row in open_opportunities)
    q4 = [
        row
        for row in open_opportunities
        if row["Design Win Date Parsed"]
        and date(2026, 10, 1) <= row["Design Win Date Parsed"] <= date(2026, 12, 31)
    ]
    support_categories = Counter(issue[3] for issue in SUPPORT_ISSUES)
    support_issue_text = "\n".join(
        (
            f"{index}. **{issue}** — {customer}; latest {latest.strftime('%d %b %Y')}; "
            f"{category}; **{status}**. {outcome}"
        )
        for index, (issue, customer, latest, category, outcome, status) in enumerate(
            SUPPORT_ISSUES, start=1
        )
    )
    reference_text = "\n".join(
        f"- **{reference_id}** — {reference_title}: {url}"
        for reference_id, reference_title, url in REFERENCES
    )
    text = f"""# Sahil FAE QBR — presenter notes

Review: {REVIEW_DATE.strftime("%d %B %Y")}, 12:00–13:00  
Prepared: {AS_OF.strftime("%d %B %Y")}  
Sources: `Sahil Report.xlsx`, `Sahil-Open Pipeline-2026.xlsx`,
`Sahil-Open Pipeline-2027.xlsx`

## Opening message

This report attributes records owned by Sahil Patni and Kasturi Rangan to Sahil. The uploaded exports
contain {len(open_opportunities)} qualified opportunities:
{fmt_value(qualified_total)} peak and {fmt_value(weighted)} weighted. The master report contains another
{len(outside_open_opportunities)} plays worth
{fmt_value(sum(row["Peak Value"] for row in outside_open_opportunities))} outside both exports.
Target attainment still cannot be calculated because target and actual-achievement data are absent.

The deck is strong enough to expose the issues, but the execution system is not yet ready to scale.
Only Sahil is named; no Arrow or Macnica specialist is committed, no distributor-linked opportunity is
present in the source, and solution/demo readiness has not been verified.

## Facts to land

- The 2026 open export contains {len(open_2026)} reporting-scope opportunities worth
  {fmt_value(sum(row["Peak Value"] for row in open_2026))}.
- The 2027 open export contains {len(open_2027)} reporting-scope opportunities worth
  {fmt_value(sum(row["Peak Value"] for row in open_2027))}.
- Combined qualified open pipeline is {fmt_value(qualified_total)} across
  {len({row["Account Name"] for row in open_opportunities})} direct customers.
- {fmt_value(sum(row["Peak Value"] for row in identify_opportunities))} remains Identify at 0% outside
  both open exports; one additional Define-stage Data Recorder worth
  {fmt_value(sum(row["Peak Value"] for row in other_outside_opportunities))} is also outside.
- Total discovered reporting-scope portfolio is {fmt_value(total)} across {len(opportunities)} opportunities.
- Outdu AI contributes {fmt_value(4_125_000)}, or {4_125_000 / qualified_total:.1%} of qualified open value.
- {len(q4)} Q4 2026 DWIN-dated plays total {fmt_value(sum(row["Peak Value"] for row in q4))};
  their weighted value is {fmt_value(sum(row["Weighted Value"] for row in q4))}.
- A genuine Top 15 can now be shown from {len(open_opportunities)} qualified open opportunities.
- Every qualified reporting-scope opportunity is marked `Altera Opportunity`; distributor account is blank.
- Treat “2026” and “2027” as the uploaded planning-bucket labels. They are not synonymous with Design
  Win Date year; the 2027 export contains some opportunities with 2026 DWIN dates.

## Account ownership and support execution

- **Primary accounts:** {", ".join(PRIMARY_ACCOUNTS)}.
- **Other strategic/supported accounts:** {", ".join(OTHER_STRATEGIC_ACCOUNTS)}.
- Philips appeared in both supplied account lists and is treated once as a primary account.
- The allocation roster contains {len(ALLOCATION_ROSTER)} customer/location rows, all naming Sahil as
  New FAE. The image spelling “Outdo” is normalized to Outdu.
- From 15 May to 28 July 2026, {len(SUPPORT_ISSUES)} distinct substantive external issue threads were
  handled: {", ".join(f"{category} {count}" for category, count in support_categories.most_common())}.
- Only {sum(issue[5] == "Resolved" for issue in SUPPORT_ISSUES)} thread explicitly records confirmed
  resolution. The remaining summaries show guidance, workaround, escalation, investigation or pending
  validation and should not be presented as closed without confirmation.
- Four of five primary accounts—Ciena, Philips, GE HealthCare and Outdu—had no substantive issue thread
  in this window. This is a proactive-engagement gap, not evidence that no technical needs exist.

### External support issue ledger

{support_issue_text}

## Altera solution story

Lead with complete, demonstrable solution paths:

- **Robotics/control:** ROS Consolidated Robot Controller, Drive-on-Chip, Drive-on-Chip with PLC,
  3×2.5G TSN example, Sensor Fusion Platform and the functional-safety flow.
- **Camera/AI:** Holoscan Sensor Bridge (MIPI to 10GbE), 4Kp60 Multi-Sensor HDR Camera,
  4Kp30 Multi-Sensor Camera with AI, Smart Camera Demo Kit, Video and Vision Processing Suite,
  FPGA AI Suite and MIPI CSI-2.
- **Partner platforms:** Arrow Eagle Board; Macnica Sulphur Agilex 5 kit and MEP100 ST2110
  SmartNIC; Critical Link vision modules and other ecosystem IP.

Be precise: NVIDIA owns Holoscan technology; Altera provides FPGA integration/reference designs.
Also verify access, release maturity, kit availability and certification scope before promising a
solution to a customer.

## Joint Altera–Arrow–Macnica model

This cannot be a Sahil-only program.

- **Altera — Sahil:** program and solution lead; architecture, benchmark method, segment message and
  specialist-DFAE escalation.
- **Arrow — name by 05 August:** at least one robotics/AI/camera-capable technical specialist;
  Arrow account map, kits/samples, workshops and opportunity follow-up.
- **Macnica — name by 05 August:** at least one robotics/AI/camera-capable technical specialist;
  Macnica account map, Sulphur/MEP100/video ecosystem, evaluations and integration support.
- **Readiness by 14 August:** at least one working robotics demo path and one working camera/AI demo
  path, with owners and repeatable benchmark instructions.
- **Targeting by 31 August:** 30 named accounts—15 Arrow and 15 Macnica—with sponsor, use case,
  installed competition and next action.
- **Proposed Q4 scorecard:** 12 workshops, 8 evaluations, 6 qualified opportunities,
  3 Develop/Design-stage plays and 2 DWINs. These targets need explicit review approval.

## Market review talk track

### Robotics

- IFR recorded 542,000 global industrial-robot installations in 2024 and 4.664 million robots in
  operation. India installed 9,100 units, up 7%, ranking sixth globally; automotive represented 45%.
- AMD's strongest practical counter is not a single device specification: it is the KR260/K26 SOM
  path with native ROS 2 and the Kria Robotics Stack.
- The Altera response should be a measured sense-to-act demonstration combining deterministic ROS 2,
  TSN, motion/control, safety and sensor fusion—not a generic FPGA feature presentation.

### Video and vision

- Interact Analysis reported a 3.9% global machine-vision decline in 2024 and forecast 1.5% growth to
  $5.7 billion in 2025. Area-scan cameras were under pressure, including competition from APAC vendors.
- AMD can lead with KV260 ease of evaluation, Vitis libraries and Versal AI Engine scale.
- Altera should prove the complete ingest→ISP→AI→output pipeline using the customer's model and sensor:
  latency distribution, power, image quality, resources, BOM and engineering effort.

### Industrial

- Rockwell's APAC survey found 94% had invested or planned to invest in AI; quality control and process
  optimization were leading use cases, while cyber standards were broadly important.
- TSN and functional safety are credible capabilities for both AMD and Altera. Avoid claiming exclusivity.
- The specific Altera proof point is Agilex 5 SoC's three hardened 2.5G TSN MACs combined with
  Drive-on-Chip and the functional-safety methodology. Quantify system consolidation and certification work.

## Competitive positioning rules

1. Do not compare Altera logic elements directly with AMD system logic cells.
2. Do not compare headline TOPS without matching model, precision, sparsity, clocks, batch and power.
3. Treat vendor power/performance claims as vendor-stated until reproduced.
4. Capture the installed AMD device, board, tool version and workload before proposing migration.
5. Win with a reproducible customer benchmark and a named de-risk plan.

## Inputs required before presenting

1. 2026 annual target.
2. YTD actual achievement and the exact definition used (revenue, bookings, DWINs, or another measure).
3. Achieved 2026 DWIN count/value.
4. Support activities completed and open technical gaps.
5. Market-share baseline by key customer.
6. Distributor account targets and joint plans.
7. Named Sales, FAE, and DFAE owners for stage exits.

## Suggested close

“I will manage the portfolio by evidence, not activity: qualify the two 0% plays, close the Q4 stage
gates, create the fifteenth strategic opportunity, and run one repeatable 2027 demand-creation motion
for robotics/control, video/vision, and industrial platforms.”

## Important definitions

- Values are shown in source units because the workbook does not identify a currency.
- “DWIN-dated” means the scheduled Design Win Date falls in that year; it does not mean the design win
  has already been achieved.
- Proposed actions in the deck are planning recommendations inferred from stage/date/product data and
  should be validated with Sales and the customer.

## References

Accessed 28 July 2026.

{reference_text}
"""
    output.parent.mkdir(parents=True, exist_ok=True)
    output.write_text(text, encoding="utf-8")


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--source", type=Path, default=Path("qbr/Sahil Report.xlsx"))
    parser.add_argument(
        "--open-2026",
        type=Path,
        default=Path("qbr/Sahil-Open Pipeline-2026.xlsx"),
    )
    parser.add_argument(
        "--open-2027",
        type=Path,
        default=Path("qbr/Sahil-Open Pipeline-2027.xlsx"),
    )
    parser.add_argument("--output-dir", type=Path, default=Path("qbr/output"))
    args = parser.parse_args()

    line_items, opportunities = read_pipeline(args.source)
    if not opportunities:
        raise SystemExit(f"No records found for reporting owners: {sorted(REPORTING_OWNERS)}")
    open_buckets = read_open_pipeline_buckets(
        {"2026": args.open_2026, "2027": args.open_2027},
        opportunities,
    )

    build_presentation(
        opportunities,
        open_buckets,
        args.output_dir / "Sahil_QBR_2026-07-29.pptx",
    )
    build_summary_workbook(
        line_items,
        opportunities,
        open_buckets,
        args.output_dir / "Sahil_QBR_Pipeline_Summary.xlsx",
    )
    build_notes(
        opportunities,
        open_buckets,
        args.output_dir / "Sahil_QBR_Presenter_Notes.md",
    )
    print(f"Generated QBR artifacts in {args.output_dir}")


if __name__ == "__main__":
    main()
